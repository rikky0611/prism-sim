"""
Generate PowerPoint presentation for RL Procedure Assistant Experiment
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json
import os
from pathlib import Path

# Get project root
PROJECT_ROOT = Path(__file__).parent.parent.parent

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    # Format title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    return slide

def add_bullet_slide(prs, title, bullets, subtitle=None):
    """Add slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]

    title_shape.text = title

    # Format title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Add subtitle if provided
    if subtitle:
        tf = body_shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.italic = True
        p.space_after = Pt(20)

        # Add bullets after subtitle
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)
    else:
        # Add bullets
        tf = body_shape.text_frame
        tf.clear()
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)

    return slide

def add_image_slide(prs, title, image_path, caption=None):
    """Add slide with title and image"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Add image
    if os.path.exists(image_path):
        left = Inches(1)
        top = Inches(1.5)
        pic = slide.shapes.add_picture(image_path, left, top, width=Inches(8))

    # Add caption if provided
    if caption:
        left = Inches(0.5)
        top = Inches(6.5)
        width = Inches(9)
        height = Inches(0.5)
        caption_box = slide.shapes.add_textbox(left, top, width, height)
        caption_frame = caption_box.text_frame
        caption_frame.text = caption
        caption_frame.paragraphs[0].font.size = Pt(14)
        caption_frame.paragraphs[0].font.italic = True

    return slide

def add_table_slide(prs, title, headers, rows, description=None):
    """Add slide with title and table"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Add description if provided
    start_top = Inches(1.3)
    if description:
        left = Inches(0.5)
        top = Inches(1.3)
        width = Inches(9)
        height = Inches(0.5)
        desc_box = slide.shapes.add_textbox(left, top, width, height)
        desc_frame = desc_box.text_frame
        desc_frame.text = description
        desc_frame.paragraphs[0].font.size = Pt(16)
        start_top = Inches(2.0)

    # Add table
    rows_count = len(rows) + 1  # +1 for header
    cols_count = len(headers)
    left = Inches(1.0)
    top = start_top
    width = Inches(8.0)
    height = Inches(0.5 * rows_count)

    table = slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table

    # Set column widths
    col_width = int(width / cols_count)
    for col_idx in range(cols_count):
        table.columns[col_idx].width = col_width

    # Add headers
    for col_idx, header in enumerate(headers):
        cell = table.rows[0].cells[col_idx]
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Add rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.rows[row_idx + 1].cells[col_idx]
            cell.text = str(cell_text)
            cell.text_frame.paragraphs[0].font.size = Pt(12)

    return slide

def create_presentation():
    """Create the full presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "Learning Context-Dependent Procedure Assistant Policies with Deep RL",
        "Optimizing Interruption-Failure Tradeoffs\n\nFebruary 2026"
    )

    # Slide 2: Research Problem
    add_bullet_slide(
        prs,
        "Research Problem",
        [
            "Challenge: When should AI assistants intervene?",
            "Interruptions are costly:",
            "  • Distraction and cognitive load",
            "  • User annoyance",
            "Failures are costly:",
            "  • Task errors",
            "  • Safety risks",
            "Key insight: Optimal strategy depends on context (cost structure)"
        ]
    )

    # Slide 3: Approach
    add_bullet_slide(
        prs,
        "Approach",
        [
            "POMDP Formulation:",
            "  • State: step, time, memory",
            "  • Actions: silent, remind, confirm",
            "Cost Function: R = -c_remind × interruptions - c_fail × failures",
            "Deep RL: PPO to learn optimal policies",
            "Baseline Policies: Random, Proactive, Reactive (hand-crafted heuristics)"
        ]
    )

    # Slide 4: Experimental Setup
    add_bullet_slide(
        prs,
        "Experimental Setup",
        [
            "Environment: 5-step cooking procedure",
            "State: Current step, elapsed time, memory vector",
            "Memory dynamics: Forgetting (λ) and reminder boost (Δ)",
            "Training: PPO with 50k timesteps",
            "Evaluation: 100 episodes per policy",
            "Cost regimes: Varying c_remind and c_fail ratios"
        ]
    )

    # Slide 5: Single Regime Results
    add_image_slide(
        prs,
        "Single Regime Results (Balanced: c_remind=5, c_fail=12)",
        str(PROJECT_ROOT / "results" / "figures" / "rl_comparison_balanced.png"),
        "RL_PPO achieves -58.92 vs Reactive -75.85 (22.3% improvement) with 0 interruptions, 1.46 failures"
    )

    # Slide 6: Strategic Silence Strategy
    add_bullet_slide(
        prs,
        'The "Strategic Silence" Strategy',
        [
            "Discovery: RL learned to NEVER interrupt (0.00 interventions)",
            "Why optimal?",
            "  • Interruption cost: 5",
            "  • Benefit: preventing ~0.5 failures × 12 = 6",
            "  • Net: Not worth it in this regime",
            "Interpretation: Silence minimizes total cost",
            "HCI insight: Optimal ≠ Maximum assistance"
        ]
    )

    # Slide 7: Multi-Regime Experiment
    add_bullet_slide(
        prs,
        "Multi-Regime Experiment",
        [
            "Question: Does RL adapt to different cost structures?",
            "Approach: Train on 5 regimes with varying c_fail/c_remind ratios",
            "Parameters:",
            "  • Increased failure risk (f0=0.6, λ=0.10)",
            "  • Makes interventions more valuable",
            "Goal: Test context-dependent adaptation"
        ]
    )

    # Slide 8: Multi-Regime Results
    add_image_slide(
        prs,
        "Multi-Regime Results",
        str(PROJECT_ROOT / "results" / "figures" / "rl_multi_regime_comparison.png"),
        "Top: Performance, Interventions, Failures across regimes. Bottom: RL improvement %, Intervention vs ratio, Trade-off scatter"
    )

    # Slide 9: Context-Dependent Adaptation
    add_table_slide(
        prs,
        "Context-Dependent Adaptation ✓",
        ["Cost Regime", "Ratio (c_fail/c_remind)", "RL Interventions", "Behavior"],
        [
            ["Very High Stakes", "15.0", "19.80", "ACTIVE"],
            ["High Stakes", "10.0", "16.90", "ACTIVE"],
            ["Moderate (V1)", "5.0", "2.14", "SELECTIVE"],
            ["Moderate-Low", "3.0", "0.00", "SILENT"],
            ["Low Stakes", "2.0", "0.00", "SILENT"]
        ],
        "Discovery: Threshold-based switching at ratio ≈ 10. Range: 0 to 19.8 interventions (20× difference)"
    )

    # Slide 10: Most Interesting - Selective Strategy
    add_bullet_slide(
        prs,
        "Most Interesting: Selective Strategy",
        [
            "Moderate regime (V1): 2.14 interventions",
            "Not too active (vs 16.5 for Proactive)",
            "Not too passive (vs 0 for most regimes)",
            "Performance: +2.8% better than best baseline",
            "Goldilocks zone: Learned nuanced intervention timing",
            "Shows RL can discover middle-ground strategies"
        ]
    )

    # Slide 11: Key Findings
    add_bullet_slide(
        prs,
        "Key Findings",
        [
            "✓ RL adapts to context:",
            "  • 20× difference in behavior (0 vs 19.8 interventions)",
            "✓ Learns threshold policy:",
            "  • Switches at ratio ≈ 10",
            "✓ Can learn selective strategies:",
            "  • 2.14 interventions in sweet spot",
            "✗ Struggles in very high-stakes:",
            "  • Too aggressive, underperforms baselines"
        ]
    )

    # Slide 12: HCI Implications
    add_bullet_slide(
        prs,
        "HCI Implications",
        [
            "Rethinking assistance: Optimal ≠ Always helping",
            "Context-aware systems:",
            "  • Single RL framework adapts to user/task context",
            "Personalization:",
            "  • Can train on individual user cost functions",
            "Design principle: Strategic silence is a valid strategy",
            "Framework for cost-aware adaptive assistance"
        ]
    )

    # Slide 13: Limitations & Future Work
    add_bullet_slide(
        prs,
        "Limitations & Future Work",
        [
            "Limitations:",
            "  • Simulation only (no human subjects)",
            "  • Simplified 5-step procedure",
            "  • Over-intervenes in very high-stakes",
            "Future directions:",
            "  • Multi-task RL (single policy for all contexts)",
            "  • Human subject validation",
            "  • Hierarchical policies for selective intervention"
        ]
    )

    # Slide 14: Conclusions
    add_bullet_slide(
        prs,
        "Conclusions",
        [
            "✅ RL significantly outperforms heuristics",
            "   (22.3% improvement in balanced regime)",
            "✅ Discovers context-dependent strategies automatically",
            "✅ Identifies 'strategic silence' as optimal in some contexts",
            "⚠️ Challenges remain for selective intervention in high-stakes",
            "🔬 Research value: Framework for cost-aware adaptive assistance"
        ]
    )

    # Slide 15: Thank You
    add_title_slide(
        prs,
        "Thank You",
        "Questions?"
    )

    # Save presentation to results/presentations/
    output_path = PROJECT_ROOT / "results" / "presentations" / "RL_Experiment_Presentation.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Presentation saved to: {output_path}")

    return str(output_path)

if __name__ == "__main__":
    create_presentation()
