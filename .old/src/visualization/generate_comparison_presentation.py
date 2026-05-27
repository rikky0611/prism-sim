"""
Generate Comparison PowerPoint: Old Model vs New Timing-Dependent Model

Compares results from weak reminders (45% effectiveness) vs strong timing-dependent
reminders (90-100% effectiveness when well-timed).
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Add project root to path
PROJECT_ROOT = Path(__file__).parent.parent.parent


def load_results(filename):
    """Load evaluation results from JSON."""
    results_path = PROJECT_ROOT / "data" / "results" / filename
    with open(results_path, 'r') as f:
        return json.load(f)


def add_title_slide(prs, title, subtitle):
    """Add title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    # Format title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)


def add_content_slide(prs, title, content_points):
    """Add content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True

    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()

    for point in content_points:
        p = text_frame.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(12)


def add_image_slide(prs, title, image_path, caption=None):
    """Add slide with image."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add image
    img_left = Inches(1)
    img_top = Inches(1.3)
    img_width = Inches(8)

    if Path(image_path).exists():
        slide.shapes.add_picture(str(image_path), img_left, img_top, width=img_width)

    # Add caption if provided
    if caption:
        cap_top = Inches(6.8)
        cap_height = Inches(0.5)
        cap_box = slide.shapes.add_textbox(left, cap_top, width, cap_height)
        cap_frame = cap_box.text_frame
        cap_frame.text = caption
        cap_frame.paragraphs[0].font.size = Pt(14)
        cap_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def create_comparison_presentation():
    """Create comprehensive comparison PowerPoint."""

    print("Creating comparison PowerPoint presentation...")

    # Load new results (timing-dependent model)
    results_new = load_results("cross_task_multi_regime_evaluation_NEW.json")
    aggregate_new = results_new['aggregate_stats']

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "Comparing Reminder Effectiveness Models",
        "Old Model (Weak 45%) vs New Model (Strong 90-100% with Timing)\\n\\n" +
        "Cross-Task Multi-Regime RL Experiments\\n" +
        "2026-02-16"
    )

    # Slide 2: Executive Summary - OLD MODEL (from documentation)
    add_content_slide(
        prs,
        "OLD MODEL Results (Weak 45% Prevention)",
        [
            "✓ Trained 21 RL models (7 tasks × 3 cost regimes)",
            f"✓ 24.21% mean improvement over baselines",
            f"✓ 100% success rate - all 21 models beat baselines",
            "✓ Strategic silence emerges across all regimes",
            "❌ PROBLEM: Very high stakes → 0.00 interruptions (counter-intuitive)",
            "❌ ROOT CAUSE: Reminders only 45% effective, silence becomes optimal"
        ]
    )

    # Slide 3: Executive Summary - NEW MODEL
    add_content_slide(
        prs,
        "NEW MODEL Results (Strong 90-100% Prevention)",
        [
            "✓ Re-trained same 21 RL models with timing-dependent effectiveness",
            f"✓ {aggregate_new['overall']['mean_improvement_pct']:.1f}% mean improvement over baselines",
            f"⚠ {aggregate_new['overall']['success_rate']*100:.0f}% success rate - " +
            f"{aggregate_new['overall']['n_models_improved']}/21 models beat baselines",
            "⚠ 3 models performed worse (make_stencil tasks)",
            "❓ QUESTION: Why did overall performance decrease?",
            "💡 HYPOTHESIS: RL hasn't learned optimal timing strategies yet"
        ]
    )

    # Slide 4: Key Difference in Models
    add_content_slide(
        prs,
        "What Changed Between Models?",
        [
            "OLD MODEL: Flat effectiveness (no timing awareness)",
            "  • Reminder 50 ticks early = same as 1 tick before",
            "  • Maximum prevention: ~45-55%",
            "  • f(m) = f0_base × exp(-k × m)",
            "",
            "NEW MODEL: Timing-dependent effectiveness",
            "  • Recency factor: r = exp(-λ × ticks_since_reminder)",
            "  • Just reminded (0-2 ticks): 97-99% prevention",
            "  • Old reminder (20+ ticks): 45-55% prevention (degrades to baseline)",
            "  • f(m,r) = f0_base × exp(-k × m) × (1 - 0.95 × r)"
        ]
    )

    # Slide 5: Overall Performance Comparison
    overall_new = aggregate_new['overall']
    add_content_slide(
        prs,
        "Overall Performance: OLD vs NEW",
        [
            "OLD MODEL (Weak 45%):",
            "  • Mean improvement: 24.21%",
            "  • Success rate: 100% (21/21 models)",
            "  • Best case: +35.77% (make_cereal)",
            "  • Worst case: +4.58% (latte_making)",
            "",
            "NEW MODEL (Strong 90-100% with Timing):",
            f"  • Mean improvement: {overall_new['mean_improvement_pct']:.2f}%",
            f"  • Success rate: {overall_new['success_rate']*100:.0f}% ({overall_new['n_models_improved']}/21)",
            f"  • Best case: +{overall_new['best_case']['improvement_pct']:.2f}% ({overall_new['best_case']['task']})",
            f"  • Worst case: {overall_new['worst_case']['improvement_pct']:.2f}% ({overall_new['worst_case']['task']})"
        ]
    )

    # Slide 6: Results by Regime - NEW MODEL
    regime_data = []
    for regime, stats in aggregate_new['by_regime'].items():
        regime_label = regime.replace('_', ' ').title()
        regime_data.append(f"{regime_label}:")
        regime_data.append(f"  • Mean improvement: {stats['mean_improvement_pct']:.2f}%")
        regime_data.append(f"  • Mean interventions: {stats['mean_rl_interruptions']:.2f}")
        regime_data.append(f"  • Success rate: {stats['success_rate']*100:.0f}%")
        regime_data.append("")

    add_content_slide(prs, "Performance by Cost Regime (NEW MODEL)", regime_data)

    # Slide 7: Interruption Pattern Comparison
    add_content_slide(
        prs,
        "Interruption Patterns: OLD vs NEW",
        [
            "OLD MODEL (Weak 45%):",
            "  • Very High Stakes (ratio=15): 0.00 interruptions",
            "  • Balanced (ratio=3): 0.47 interruptions",
            "  • Moderate Low (ratio=2): 3.37 interruptions",
            "  → PATTERN: Lower stakes → MORE interventions (counter-intuitive!)",
            "",
            "NEW MODEL (Strong 90-100%):",
            f"  • Very High Stakes (ratio=15): {aggregate_new['by_regime']['very_high_stakes']['mean_rl_interruptions']:.2f} interruptions",
            f"  • Balanced (ratio=3): {aggregate_new['by_regime']['balanced']['mean_rl_interruptions']:.2f} interruptions",
            f"  • Moderate Low (ratio=2): {aggregate_new['by_regime']['moderate_low']['mean_rl_interruptions']:.2f} interruptions",
            "  → PATTERN: Still similar (silence still dominates)"
        ]
    )

    # Slide 8: Why Did Performance Decrease?
    add_content_slide(
        prs,
        "Why Did Overall Performance Decrease?",
        [
            "Hypothesis 1: Timing Learning Difficulty",
            "  • New model requires learning WHEN to remind (timing strategy)",
            "  • 50k timesteps may be insufficient for timing optimization",
            "  • Old model only needed to learn IF to remind (simpler)",
            "",
            "Hypothesis 2: Exploration Challenges",
            "  • Fast recency decay (λ=0.20, half-life ~3.5 ticks) is very strict",
            "  • RL needs to discover narrow timing windows (0-5 ticks before step)",
            "  • Random exploration unlikely to hit these windows often",
            "",
            "Hypothesis 3: Credit Assignment",
            "  • Reward signal comes at step completion, delayed from reminder",
            "  • Hard to attribute success to timing of reminder vs just giving reminder"
        ]
    )

    # Slide 9: Performance Heatmap
    fig_path = PROJECT_ROOT / "results" / "figures" / "cross_task_performance_heatmap.png"
    add_image_slide(
        prs,
        "Performance Heatmap (NEW MODEL): Task × Regime",
        fig_path,
        "RL mean reward and improvement % for all 21 task-regime combinations"
    )

    # Slide 10: Interruption Patterns
    fig_path = PROJECT_ROOT / "results" / "figures" / "cross_task_interruptions_heatmap.png"
    add_image_slide(
        prs,
        "Intervention Frequency (NEW MODEL)",
        fig_path,
        "RL learns similar silence patterns despite stronger reminder effectiveness"
    )

    # Slide 11: Which Models Got Worse?
    add_content_slide(
        prs,
        "Models That Performed Worse (NEW MODEL)",
        [
            "3 models showed negative improvement:",
            "",
            "1. make_stencil / very_high_stakes: -22.65%",
            "   • Complex task (17 steps, safety-critical)",
            "   • High criticality multipliers (up to 2.5×)",
            "   • Needs precise timing for many critical steps",
            "",
            "2. make_stencil / balanced: -14.54%",
            "3. make_stencil / moderate_low: -14.54%",
            "",
            "PATTERN: All failures are the same complex task (make_stencil)",
            "→ Suggests timing learning is especially hard for complex procedures"
        ]
    )

    # Slide 12: What We Learned
    add_content_slide(
        prs,
        "Key Insights from Model Comparison",
        [
            "1. Stronger reminders ≠ automatically better RL performance",
            "   • Need to learn HOW to use the stronger tool effectively",
            "",
            "2. Timing adds complexity to the learning problem",
            "   • Old model: Binary decision (remind or not)",
            "   • New model: Continuous timing decision (when to remind)",
            "",
            "3. Simple flat effectiveness has an advantage:",
            "   • No timing pressure - reminder value doesn't decay rapidly",
            "   • Easier credit assignment - just learn which steps need reminders",
            "",
            "4. Strategic silence still emerges:",
            "   • Even with 90-100% prevention, interruption costs dominate",
            "   • High baseline failure rate (f0=0.6) means many steps fail anyway"
        ]
    )

    # Slide 13: Implications for Real-World Assistants
    add_content_slide(
        prs,
        "Implications for Real-World Systems",
        [
            "Finding 1: Perfect effectiveness alone doesn't guarantee better assistance",
            "  • System must also learn optimal TIMING of interventions",
            "  • May require more sophisticated training (curriculum learning, demonstrations)",
            "",
            "Finding 2: Simpler models can outperform complex ones",
            "  • Flat effectiveness (old model) easier to optimize than timing-dependent (new)",
            "  • Trade-off: Realism vs. Learnability",
            "",
            "Finding 3: Strategic silence is robust",
            "  • Emerges even when reminders are highly effective",
            "  • Interruption costs matter more than reminder strength",
            "",
            "Design Recommendation:",
            "  • Start with simpler effectiveness models for RL training",
            "  • Add timing complexity only after mastering basic intervention strategy"
        ]
    )

    # Slide 14: Future Directions
    add_content_slide(
        prs,
        "Next Steps to Improve NEW MODEL",
        [
            "1. Longer training (200k timesteps)",
            "   • Give RL more time to discover timing strategies",
            "",
            "2. Curriculum learning",
            "   • Start with flat effectiveness, gradually add timing constraints",
            "   • Learn intervention first, then refine timing",
            "",
            "3. Adjust recency decay rate",
            "   • Try slower decay (λ=0.10 instead of 0.20)",
            "   • Wider effectiveness window (10-20 ticks instead of 5-10)",
            "",
            "4. Hierarchical RL",
            "   • High-level policy: decide which steps to remind",
            "   • Low-level policy: decide when to give reminder",
            "",
            "5. Expert demonstrations",
            "   • Show optimal timing patterns to bootstrap learning"
        ]
    )

    # Slide 15: Research Questions
    add_content_slide(
        prs,
        "Open Research Questions",
        [
            "Q1: What's the optimal recency decay rate for RL learning?",
            "   • Fast (λ=0.20): Realistic but hard to learn",
            "   • Slow (λ=0.05): Easier to learn but unrealistic",
            "",
            "Q2: Can we design reward shaping to help timing discovery?",
            "   • Bonus for well-timed reminders (just before step completion)?",
            "",
            "Q3: Does timing matter more for complex vs simple tasks?",
            "   • make_stencil (17 steps) failed, but make_cereal (8 steps) succeeded",
            "",
            "Q4: Would human demonstrations accelerate timing learning?",
            "   • Imitation learning + RL fine-tuning?",
            "",
            "Q5: Is there a sweet spot between flat and timing-dependent?",
            "   • Medium decay rate that balances realism and learnability?"
        ]
    )

    # Slide 16: Conclusions
    add_content_slide(
        prs,
        "Conclusions: Comparing the Two Models",
        [
            "OLD MODEL (Weak 45%, Flat Effectiveness):",
            "  ✓ 24.21% mean improvement, 100% success rate",
            "  ✓ Easier for RL to learn (no timing pressure)",
            "  ❌ Unrealistic (reminders don't degrade over time)",
            "  ❌ Led to counter-intuitive silence in high stakes",
            "",
            "NEW MODEL (Strong 90-100%, Timing-Dependent):",
            "  ✓ 13.92% mean improvement, 85.7% success rate",
            "  ✓ More realistic (reminders degrade naturally)",
            "  ❌ Harder for RL to learn (timing complexity)",
            "  ❌ Some complex tasks performed worse",
            "",
            "Key Takeaway: Realism vs. Learnability tradeoff",
            "  • More realistic models aren't always better for RL optimization",
            "  • May need specialized training approaches for timing-dependent systems"
        ]
    )

    # Slide 17: Recommendations
    add_content_slide(
        prs,
        "Recommendations for Future Work",
        [
            "For RL Training:",
            "  1. Use curriculum learning (flat → timing-dependent)",
            "  2. Increase training duration (200k+ timesteps)",
            "  3. Try slower recency decay (λ=0.10 vs 0.20)",
            "",
            "For Model Design:",
            "  4. Balance realism with learnability",
            "  5. Consider hybrid: flat for learning, timing for deployment",
            "",
            "For Research:",
            "  6. Investigate reward shaping for timing discovery",
            "  7. Compare hierarchical vs end-to-end timing learning",
            "  8. Test human demonstration bootstrapping",
            "",
            "For Practice:",
            "  9. Start simple: master intervention strategy before adding timing",
            "  10. Monitor complex tasks separately (they need special treatment)"
        ]
    )

    # Save presentation
    output_path = PROJECT_ROOT / "results" / "presentations" / "Model_Comparison_Old_vs_New.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    print(f"✓ Presentation saved to: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")

    return output_path


def main():
    """Main entry point."""
    print("="*70)
    print("GENERATING MODEL COMPARISON PRESENTATION")
    print("="*70)
    print()

    output_path = create_comparison_presentation()

    print()
    print("="*70)
    print("PRESENTATION COMPLETE")
    print("="*70)
    print(f"PowerPoint file: {output_path}")
    print(f"Size: {output_path.stat().st_size / 1024:.1f} KB")
    print("="*70)


if __name__ == '__main__':
    main()
