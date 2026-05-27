"""
Generate Detailed NEW MODEL Presentation

Creates a comprehensive PowerPoint presentation focused on the NEW MODEL
(timing-dependent 90-100% reminder effectiveness) with detailed task-by-task
analysis including:
- Task descriptions and step-by-step information
- Per-step failure costs with criticality multipliers
- Observation noise details
- Performance comparison with baselines
- Intervention frequency and timing patterns

Author: Claude Sonnet 4.5
Date: 2026-02-16
"""

import json
from pathlib import Path
from typing import Dict, List, Any, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Add parent directories to path for imports
import sys
PROJECT_ROOT = Path(__file__).parent.parent.parent
sys.path.insert(0, str(PROJECT_ROOT / "src" / "experiments"))

from task_definitions import load_task_definitions, TaskDefinition


# =============================================================================
# CONFIGURATION
# =============================================================================

TITLE_FONT_SIZE = Pt(32)
SUBTITLE_FONT_SIZE = Pt(18)
HEADING_FONT_SIZE = Pt(24)
BODY_FONT_SIZE = Pt(14)
SMALL_FONT_SIZE = Pt(11)

COLOR_TITLE = RGBColor(31, 78, 120)
COLOR_HEADING = RGBColor(68, 114, 196)
COLOR_SUCCESS = RGBColor(0, 128, 0)
COLOR_WARNING = RGBColor(255, 140, 0)
COLOR_FAILURE = RGBColor(192, 0, 0)
COLOR_NEUTRAL = RGBColor(100, 100, 100)


# =============================================================================
# UTILITY FUNCTIONS
# =============================================================================

def set_text_properties(text_frame, text: str, font_size: Pt,
                        bold: bool = False, color: RGBColor = None,
                        alignment: PP_ALIGN = None):
    """Set text properties for a text frame."""
    text_frame.text = text
    p = text_frame.paragraphs[0]
    p.font.size = font_size
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    if alignment:
        p.alignment = alignment


def add_title_slide(prs: Presentation, title: str, subtitle: str = ""):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    set_text_properties(title_box.text_frame, title, TITLE_FONT_SIZE,
                       bold=True, color=COLOR_TITLE, alignment=PP_ALIGN.CENTER)

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.8))
        set_text_properties(subtitle_box.text_frame, subtitle, SUBTITLE_FONT_SIZE,
                           color=COLOR_HEADING, alignment=PP_ALIGN.CENTER)

    return slide


def add_section_title_slide(prs: Presentation, title: str, subtitle: str = ""):
    """Add a section title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    set_text_properties(title_box.text_frame, title, TITLE_FONT_SIZE,
                       bold=True, color=COLOR_TITLE, alignment=PP_ALIGN.CENTER)

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
        set_text_properties(subtitle_box.text_frame, subtitle, BODY_FONT_SIZE,
                           color=COLOR_HEADING, alignment=PP_ALIGN.CENTER)

    return slide


def add_content_slide(prs: Presentation, title: str) -> Any:
    """Add a content slide with title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    set_text_properties(title_box.text_frame, title, HEADING_FONT_SIZE,
                       bold=True, color=COLOR_HEADING)

    return slide


# =============================================================================
# SLIDE CREATION FUNCTIONS
# =============================================================================

def create_executive_summary_slide(prs: Presentation, results: Dict, tasks: Dict):
    """Slide 1: Executive Summary"""
    slide = add_content_slide(prs, "Executive Summary: NEW MODEL Performance")

    # Calculate overall statistics
    all_improvements = []
    success_count = 0
    total_count = 0

    for regime_name, regime_results in results.items():
        for task_name, task_result in regime_results.items():
            if 'improvement_pct' in task_result:
                improvement = task_result['improvement_pct']
                all_improvements.append(improvement)
                if improvement > 0:
                    success_count += 1
                total_count += 1

    mean_improvement = sum(all_improvements) / len(all_improvements) if all_improvements else 0
    success_rate = (success_count / total_count * 100) if total_count > 0 else 0
    best_improvement = max(all_improvements) if all_improvements else 0
    worst_improvement = min(all_improvements) if all_improvements else 0

    # Content
    y_pos = 1.2
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(8.4), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    # Model info
    p = tf.paragraphs[0]
    p.text = "Model Configuration:"
    p.font.size = BODY_FONT_SIZE
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADING

    details = [
        "• Reminder Effectiveness: 90-100% (timing-dependent, recency decay λ=0.20)",
        "• Training: 50,000 timesteps per model, PPO algorithm",
        "• Evaluation: 100 episodes per model",
        "• Tasks: 7 procedural tasks (8-20 steps)",
        "• Cost Regimes: 3 (Very High Stakes, Balanced, Moderate Low)",
        "",
        "Overall Performance:",
        f"• Mean Improvement over Baselines: {mean_improvement:.2f}%",
        f"• Success Rate: {success_rate:.1f}% ({success_count}/{total_count} models)",
        f"• Best Case: +{best_improvement:.2f}%",
        f"• Worst Case: {worst_improvement:.2f}%",
        "",
        "Key Insight:",
        "Despite 90-100% reminder effectiveness, timing-dependent model requires",
        "learning WHEN to remind (continuous timing) vs IF to remind (binary decision).",
        "50k timesteps sufficient for simple tasks, insufficient for complex tasks."
    ]

    for detail in details:
        p = tf.add_paragraph()
        p.text = detail
        p.font.size = SMALL_FONT_SIZE if detail.startswith("•") else BODY_FONT_SIZE
        p.font.bold = ":" in detail and not detail.startswith("•")
        if "Worst Case:" in detail and worst_improvement < 0:
            p.font.color.rgb = COLOR_FAILURE


def create_model_details_slide(prs: Presentation):
    """Slide 2: Model Technical Details"""
    slide = add_content_slide(prs, "NEW MODEL: Technical Specification")

    y_pos = 1.2
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(8.4), Inches(4.8))
    tf = content_box.text_frame
    tf.word_wrap = True

    # Timing-Dependent Effectiveness
    p = tf.paragraphs[0]
    p.text = "Timing-Dependent Reminder Effectiveness Model"
    p.font.size = BODY_FONT_SIZE
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADING

    details = [
        "",
        "Dual-Component Memory System:",
        "• Base Memory (m): Long-term procedural knowledge, slow decay (λ=0.05)",
        "• Recency Factor (r): Short-term reminder freshness, fast decay (λ=0.20)",
        "",
        "Combined Failure Model:",
        "f(m, r) = f0_base × exp(-k × m) × (1 - ε_recency × r)",
        "",
        "where:",
        "• f0_base = 0.6 (60% baseline failure probability)",
        "• k = 2.0 (memory effectiveness coefficient)",
        "• ε_recency = 0.95 (95% maximum prevention from recent reminder)",
        "• r = exp(-0.20 × ticks_since_reminder)",
        "",
        "Effectiveness Curve:",
        "• 0-2 ticks after reminder: 97-99% prevention ✓",
        "• 3-5 ticks after reminder: 73-97% prevention ✓",
        "• 6-10 ticks: 64-92% prevention",
        "• 20+ ticks: ~46% prevention (back to base memory only)",
        "",
        "Observation Noise:",
        "• 20% Gaussian noise on all observations (memory states, step progress)",
        "• Simulates partial observability in real-world scenarios"
    ]

    for detail in details:
        if detail:
            p = tf.add_paragraph()
            p.text = detail
            p.font.size = SMALL_FONT_SIZE if detail.startswith("•") else BODY_FONT_SIZE
            p.font.bold = ":" in detail and not detail.startswith("•")
            if detail.startswith("f(m") or detail.startswith("where"):
                p.font.color.rgb = COLOR_NEUTRAL


def create_task_overview_slide(prs: Presentation, task_def: TaskDefinition,
                                task_results: Dict, regime_name: str, regime_info: Dict):
    """Create detailed slide for a single task-regime combination."""
    slide = add_content_slide(prs, f"{task_def.task_name.replace('_', ' ').title()} - {regime_name.replace('_', ' ').title()}")

    # Left column: Task info
    left_x = 0.5
    y_pos = 1.1

    # Task description box
    desc_box = slide.shapes.add_textbox(Inches(left_x), Inches(y_pos), Inches(4.5), Inches(1.2))
    tf = desc_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Task Overview"
    p.font.size = BODY_FONT_SIZE
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADING

    task_info = [
        f"Domain: {task_def.domain.capitalize()}",
        f"Steps: {task_def.n_steps}",
        f"Base Failure Cost: {task_def.base_failure_cost}",
        f"Interruption Cost: {task_def.interruption_cost}",
    ]

    for info in task_info:
        p = tf.add_paragraph()
        p.text = f"• {info}"
        p.font.size = SMALL_FONT_SIZE

    # Cost regime info
    y_pos += 1.3
    regime_box = slide.shapes.add_textbox(Inches(left_x), Inches(y_pos), Inches(4.5), Inches(0.8))
    tf = regime_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Cost Regime"
    p.font.size = BODY_FONT_SIZE
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADING

    regime_details = [
        f"Interruption Cost: {regime_info['c_remind']}",
        f"Failure Cost: {regime_info['c_fail']}",
        f"Ratio: {regime_info['ratio']:.1f} ({regime_info['description']})"
    ]

    for detail in regime_details:
        p = tf.add_paragraph()
        p.text = f"• {detail}"
        p.font.size = SMALL_FONT_SIZE

    # Step criticalities
    y_pos += 1.0
    steps_box = slide.shapes.add_textbox(Inches(left_x), Inches(y_pos), Inches(4.5), Inches(3.0))
    tf = steps_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Step-by-Step Failure Costs"
    p.font.size = BODY_FONT_SIZE
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADING

    # Show top critical steps
    step_costs = [(i, step, task_def.get_step_failure_cost(i))
                  for i, step in enumerate(task_def.steps)]
    step_costs_sorted = sorted(step_costs, key=lambda x: x[2], reverse=True)

    p = tf.add_paragraph()
    p.text = f"Critical Steps (Top {min(8, len(step_costs))}):"
    p.font.size = SMALL_FONT_SIZE
    p.font.bold = True

    for i, (idx, step, cost) in enumerate(step_costs_sorted[:8]):
        p = tf.add_paragraph()
        p.text = f"  {idx+1}. {step.name}: {cost:.1f} (×{step.criticality:.1f})"
        p.font.size = Pt(10)
        if step.criticality >= 2.0:
            p.font.color.rgb = COLOR_FAILURE
        elif step.criticality >= 1.5:
            p.font.color.rgb = COLOR_WARNING

    # Right column: Performance results
    right_x = 5.2
    y_pos = 1.1

    # Performance comparison
    perf_box = slide.shapes.add_textbox(Inches(right_x), Inches(y_pos), Inches(4.5), Inches(2.5))
    tf = perf_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Performance Comparison"
    p.font.size = BODY_FONT_SIZE
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADING

    policies = task_results.get('policies', {})

    # Create comparison table
    p = tf.add_paragraph()
    p.text = f"{'Policy':<15} {'Reward':<10} {'Fails':<8} {'Ints':<8}"
    p.font.size = Pt(10)
    p.font.bold = True

    for policy_name in ['Random', 'Proactive', 'Reactive', 'RL_PPO']:
        if policy_name in policies:
            policy = policies[policy_name]
            p = tf.add_paragraph()
            reward = policy['mean_reward']
            failures = policy['mean_failures']
            interruptions = policy['mean_interruptions']
            p.text = f"{policy_name:<15} {reward:<10.1f} {failures:<8.2f} {interruptions:<8.2f}"
            p.font.size = Pt(10)
            if policy_name == 'RL_PPO':
                p.font.bold = True
                p.font.color.rgb = COLOR_SUCCESS if task_results.get('improvement_pct', 0) > 0 else COLOR_FAILURE

    # Improvement summary
    y_pos += 2.6
    summary_box = slide.shapes.add_textbox(Inches(right_x), Inches(y_pos), Inches(4.5), Inches(1.2))
    tf = summary_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "RL Performance"
    p.font.size = BODY_FONT_SIZE
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADING

    improvement = task_results.get('improvement_pct', 0)
    best_baseline = task_results.get('best_baseline', 'Unknown')

    summary_items = [
        f"Improvement over Best Baseline: {improvement:+.2f}%",
        f"Best Baseline: {best_baseline}",
        f"Status: {'✓ Success' if improvement > 0 else '✗ Failed' if improvement < 0 else '○ Neutral'}"
    ]

    for item in summary_items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = SMALL_FONT_SIZE
        if 'Success' in item:
            p.font.color.rgb = COLOR_SUCCESS
        elif 'Failed' in item:
            p.font.color.rgb = COLOR_FAILURE

    # Intervention timing
    y_pos += 1.3
    timing_box = slide.shapes.add_textbox(Inches(right_x), Inches(y_pos), Inches(4.5), Inches(1.8))
    tf = timing_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Intervention Patterns"
    p.font.size = BODY_FONT_SIZE
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADING

    if 'RL_PPO' in policies:
        rl_policy = policies['RL_PPO']
        interruptions = rl_policy['mean_interruptions']
        failures = rl_policy['mean_failures']

        # Calculate intervention rate
        intervention_rate = (interruptions / task_def.n_steps * 100) if task_def.n_steps > 0 else 0

        timing_items = [
            f"Mean Interruptions: {interruptions:.2f}",
            f"Intervention Rate: {intervention_rate:.1f}% of steps",
            f"Mean Failures: {failures:.2f}",
            f"Strategy: {'Proactive' if interruptions > 5 else 'Strategic Silence' if interruptions < 2 else 'Selective'}"
        ]

        for item in timing_items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = SMALL_FONT_SIZE


def create_regime_summary_slide(prs: Presentation, regime_name: str,
                                regime_info: Dict, regime_results: Dict,
                                tasks: Dict):
    """Create summary slide for a cost regime across all tasks."""
    slide = add_content_slide(prs, f"Regime Summary: {regime_name.replace('_', ' ').title()}")

    y_pos = 1.1

    # Regime description
    desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(8.4), Inches(0.8))
    tf = desc_box.text_frame

    p = tf.paragraphs[0]
    p.text = f"{regime_info['description']}"
    p.font.size = BODY_FONT_SIZE
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADING

    p = tf.add_paragraph()
    p.text = f"c_remind={regime_info['c_remind']}, c_fail={regime_info['c_fail']}, ratio={regime_info['ratio']:.1f}"
    p.font.size = SMALL_FONT_SIZE

    # Results table
    y_pos += 1.0
    table_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(8.4), Inches(4.0))
    tf = table_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Performance Across All Tasks"
    p.font.size = BODY_FONT_SIZE
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADING

    # Header
    p = tf.add_paragraph()
    p.text = f"{'Task':<18} {'Steps':<7} {'Improvement':<13} {'Interrupts':<12} {'Failures':<10}"
    p.font.size = Pt(11)
    p.font.bold = True

    # Calculate statistics
    improvements = []
    success_count = 0

    for task_name, task_result in regime_results.items():
        task_def = tasks.get(task_name)
        if not task_def:
            continue

        improvement = task_result.get('improvement_pct', 0)
        improvements.append(improvement)
        if improvement > 0:
            success_count += 1

        rl_policy = task_result.get('policies', {}).get('RL_PPO', {})
        interruptions = rl_policy.get('mean_interruptions', 0)
        failures = rl_policy.get('mean_failures', 0)

        p = tf.add_paragraph()
        task_display = task_name.replace('_', ' ')[:16]
        p.text = f"{task_display:<18} {task_def.n_steps:<7} {improvement:>+7.2f}% {'':>5} {interruptions:>6.2f} {'':>5} {failures:>5.2f}"
        p.font.size = Pt(10)

        if improvement > 0:
            p.font.color.rgb = COLOR_SUCCESS
        elif improvement < 0:
            p.font.color.rgb = COLOR_FAILURE

    # Summary statistics
    mean_improvement = sum(improvements) / len(improvements) if improvements else 0
    success_rate = (success_count / len(improvements) * 100) if improvements else 0

    p = tf.add_paragraph()
    p.text = ""  # Blank line

    p = tf.add_paragraph()
    p.text = f"Mean Improvement: {mean_improvement:+.2f}%  |  Success Rate: {success_rate:.1f}% ({success_count}/{len(improvements)} tasks)"
    p.font.size = SMALL_FONT_SIZE
    p.font.bold = True
    if mean_improvement > 0:
        p.font.color.rgb = COLOR_SUCCESS
    elif mean_improvement < 0:
        p.font.color.rgb = COLOR_FAILURE


def create_conclusions_slide(prs: Presentation, results: Dict):
    """Final slide: Key Insights and Conclusions"""
    slide = add_content_slide(prs, "Key Insights: NEW MODEL Analysis")

    y_pos = 1.2
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(8.4), Inches(4.8))
    tf = content_box.text_frame
    tf.word_wrap = True

    # Calculate overall stats
    all_improvements = []
    failed_tasks = []

    for regime_name, regime_results in results.items():
        for task_name, task_result in regime_results.items():
            improvement = task_result.get('improvement_pct', 0)
            all_improvements.append(improvement)
            if improvement < 0:
                failed_tasks.append(f"{task_name} ({regime_name})")

    mean_improvement = sum(all_improvements) / len(all_improvements) if all_improvements else 0
    success_rate = sum(1 for x in all_improvements if x > 0) / len(all_improvements) * 100

    p = tf.paragraphs[0]
    p.text = "Key Findings"
    p.font.size = BODY_FONT_SIZE
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADING

    findings = [
        f"1. Overall Performance: {mean_improvement:.2f}% mean improvement, {success_rate:.1f}% success rate",
        "",
        "2. Timing Learning Complexity:",
        "   • NEW MODEL requires learning WHEN to remind (continuous timing)",
        "   • OLD MODEL only required learning IF to remind (binary decision)",
        "   • 50k timesteps sufficient for simple tasks (8-9 steps)",
        "   • Insufficient for complex tasks (17+ steps, safety-critical)",
        "",
        "3. Task Complexity Effects:",
        f"   • Simple tasks (make_cereal, make_coffee): Success ✓",
        f"   • Complex tasks (make_stencil): Failed in all regimes ✗",
        "   • Timing discovery difficulty scales with task complexity",
        "",
        "4. Realism vs. Learnability Tradeoff:",
        "   • 90-100% reminder effectiveness is realistic",
        "   • But creates harder learning problem for RL",
        "   • Need specialized training approaches:",
        "     - Longer training (200k timesteps)",
        "     - Curriculum learning (flat → timing-dependent)",
        "     - Slower recency decay (wider effectiveness window)",
        "",
        "5. Strategic Silence Still Emerges:",
        "   • High baseline failure rate (f0=0.6) creates difficulty",
        "   • Even with strong reminders, interruption costs matter",
        "   • Context-dependent strategies discovered"
    ]

    for finding in findings:
        if finding:
            p = tf.add_paragraph()
            p.text = finding
            p.font.size = SMALL_FONT_SIZE if finding.startswith("   ") else BODY_FONT_SIZE
            p.font.bold = finding[0].isdigit() or ":" in finding[:30]
            if "Failed" in finding or "✗" in finding:
                p.font.color.rgb = COLOR_FAILURE
            elif "Success" in finding or "✓" in finding:
                p.font.color.rgb = COLOR_SUCCESS


# =============================================================================
# MAIN FUNCTION
# =============================================================================

def generate_presentation():
    """Generate the complete NEW MODEL detailed presentation."""

    print("="*80)
    print("GENERATING NEW MODEL DETAILED PRESENTATION")
    print("="*80)

    # Load data
    print("\n1. Loading data...")
    results_path = PROJECT_ROOT / "data" / "results" / "cross_task_multi_regime_evaluation_NEW.json"

    with open(results_path, 'r') as f:
        data = json.load(f)

    results = data['results']
    regimes = data['regimes']
    tasks_dict = load_task_definitions()

    print(f"   ✓ Loaded results for {len(results)} regimes")
    print(f"   ✓ Loaded {len(tasks_dict)} task definitions")

    # Create presentation
    print("\n2. Creating presentation...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    add_title_slide(prs,
                   "NEW MODEL: Detailed Analysis",
                   "Timing-Dependent Reminder Effectiveness (90-100% Prevention)")
    print("   ✓ Created title slide")

    # Executive summary
    create_executive_summary_slide(prs, results, tasks_dict)
    print("   ✓ Created executive summary")

    # Model technical details
    create_model_details_slide(prs)
    print("   ✓ Created model details slide")

    # For each regime, create summary + individual task slides
    for regime_name, regime_info in regimes.items():
        regime_results = results[regime_name]

        # Regime summary slide
        create_regime_summary_slide(prs, regime_name, regime_info, regime_results, tasks_dict)
        print(f"   ✓ Created summary for {regime_name}")

        # Individual task slides
        for task_name in sorted(regime_results.keys(),
                              key=lambda t: tasks_dict[t].n_steps):
            task_def = tasks_dict[task_name]
            task_result = regime_results[task_name]

            create_task_overview_slide(prs, task_def, task_result,
                                      regime_name, regime_info)
            print(f"      ✓ Created slide for {task_name}")

    # Conclusions slide
    create_conclusions_slide(prs, results)
    print("   ✓ Created conclusions slide")

    # Save presentation
    print("\n3. Saving presentation...")
    output_path = PROJECT_ROOT / "results" / "presentations" / "NEW_MODEL_Detailed_Analysis.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    file_size = output_path.stat().st_size / 1024
    print(f"   ✓ Saved to: {output_path}")
    print(f"   ✓ File size: {file_size:.1f} KB")
    print(f"   ✓ Total slides: {len(prs.slides)}")

    print("\n" + "="*80)
    print("PRESENTATION GENERATION COMPLETE")
    print("="*80)
    print(f"\nPresentation structure:")
    print(f"  • 1 title slide")
    print(f"  • 1 executive summary")
    print(f"  • 1 technical details slide")
    print(f"  • 3 regime summaries")
    print(f"  • {len(tasks_dict) * len(regimes)} task detail slides (7 tasks × 3 regimes)")
    print(f"  • 1 conclusions slide")
    print(f"  • Total: {len(prs.slides)} slides")

    return output_path


if __name__ == "__main__":
    output_path = generate_presentation()
    print(f"\n✓ Open the presentation at: {output_path}")
