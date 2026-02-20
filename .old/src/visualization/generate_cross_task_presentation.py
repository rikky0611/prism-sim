"""
Generate PowerPoint Presentation for Cross-Task Multi-Regime RL Results

Creates a comprehensive presentation with findings, visualizations, and insights.

Usage:
    python generate_cross_task_presentation.py
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Add project root to path
PROJECT_ROOT = Path(__file__).parent.parent.parent


def load_results():
    """Load evaluation results from JSON."""
    results_path = PROJECT_ROOT / "data" / "results" / "cross_task_multi_regime_evaluation.json"
    with open(results_path, 'r') as f:
        return json.load(f)


def add_title_slide(prs, title, subtitle):
    """Add title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    # Format title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)


def add_content_slide(prs, title, content_points):
    """Add content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True

    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()

    for point in content_points:
        p = text_frame.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(12)


def add_image_slide(prs, title, image_path, caption=None):
    """Add slide with image."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add image
    img_left = Inches(1)
    img_top = Inches(1.3)
    img_width = Inches(8)

    if Path(image_path).exists():
        slide.shapes.add_picture(str(image_path), img_left, img_top, width=img_width)

    # Add caption if provided
    if caption:
        cap_top = Inches(6.8)
        cap_height = Inches(0.5)
        cap_box = slide.shapes.add_textbox(left, cap_top, width, cap_height)
        cap_frame = cap_box.text_frame
        cap_frame.text = caption
        cap_frame.paragraphs[0].font.size = Pt(14)
        cap_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_results_table_slide(prs, title, results):
    """Add slide with results table."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Add table content as text
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()

    for key, value in results.items():
        p = text_frame.add_paragraph()
        p.text = f"{key}: {value}"
        p.font.size = Pt(16)
        p.space_after = Pt(8)


def create_presentation():
    """Create comprehensive PowerPoint presentation."""

    print("Creating PowerPoint presentation...")

    # Load results
    results = load_results()
    aggregate = results['aggregate_stats']

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "Cross-Task Multi-Regime RL Experiments",
        "Learning Context-Dependent Procedure Assistant Policies\n" +
        "Across 7 Tasks and 3 Cost Regimes\n\n" +
        "2026-02-16"
    )

    # Slide 2: Executive Summary
    add_content_slide(
        prs,
        "Executive Summary",
        [
            "✓ Trained 21 RL models (7 tasks × 3 cost regimes) in 31.9 minutes",
            f"✓ {aggregate['overall']['mean_improvement_pct']:.1f}% mean improvement over baselines",
            f"✓ 100% success rate - all {aggregate['overall']['n_models_evaluated']} models beat baselines",
            "✓ Key discovery: Strategic silence emerges across all regimes",
            "✓ Surprising finding: Lower stakes → MORE interventions",
            "✓ Simple tasks show largest gains (30-36% improvement)"
        ]
    )

    # Slide 3: Experimental Setup
    add_content_slide(
        prs,
        "Experimental Setup",
        [
            "Tasks (7): make_cereal, make_coffee, make_tea, make_sandwich,\n    cooking, make_stencil, latte_making",
            "Cost Regimes (3):",
            "  • Very High Stakes: c_int=2, c_fail=30 (ratio=15)",
            "  • Balanced: c_int=5, c_fail=15 (ratio=3)",
            "  • Moderate Low: c_int=5, c_fail=10 (ratio=2)",
            "Training: 50,000 timesteps per model using PPO",
            "Evaluation: 100 episodes vs 3 baselines (Random, Proactive, Reactive)"
        ]
    )

    # Slide 4: Overall Results
    overall = aggregate['overall']
    add_content_slide(
        prs,
        "Overall Performance Results",
        [
            f"Mean Improvement: {overall['mean_improvement_pct']:.2f}%",
            f"Median Improvement: {overall['median_improvement_pct']:.2f}%",
            f"Success Rate: {overall['success_rate']*100:.0f}% ({overall['n_models_improved']}/{overall['n_models_evaluated']} models)",
            "",
            f"Best Case: {overall['best_case']['regime'].replace('_', ' ').title()} / " +
            f"{overall['best_case']['task'].replace('_', ' ').title()}",
            f"    → {overall['best_case']['improvement_pct']:.2f}% improvement",
            "",
            f"Worst Case: {overall['worst_case']['regime'].replace('_', ' ').title()} / " +
            f"{overall['worst_case']['task'].replace('_', ' ').title()}",
            f"    → {overall['worst_case']['improvement_pct']:.2f}% improvement"
        ]
    )

    # Slide 5: Results by Regime
    regime_data = []
    for regime, stats in aggregate['by_regime'].items():
        regime_label = regime.replace('_', ' ').title()
        regime_data.append(f"{regime_label}:")
        regime_data.append(f"  • Mean improvement: {stats['mean_improvement_pct']:.2f}%")
        regime_data.append(f"  • Mean interventions: {stats['mean_rl_interruptions']:.2f}")
        regime_data.append(f"  • Success rate: {stats['success_rate']*100:.0f}%")
        regime_data.append("")

    add_content_slide(prs, "Performance by Cost Regime", regime_data)

    # Slide 6: Results by Task
    task_data = []
    for task, stats in list(aggregate['by_task'].items())[:4]:  # First 4 tasks
        task_label = task.replace('_', ' ').title()
        best_regime = stats['best_regime'].replace('_', ' ').title()
        task_data.append(f"{task_label}: {stats['mean_improvement_pct']:.2f}% (best: {best_regime})")

    task_data.append("")

    for task, stats in list(aggregate['by_task'].items())[4:]:  # Remaining tasks
        task_label = task.replace('_', ' ').title()
        best_regime = stats['best_regime'].replace('_', ' ').title()
        task_data.append(f"{task_label}: {stats['mean_improvement_pct']:.2f}% (best: {best_regime})")

    add_content_slide(prs, "Performance by Task", task_data)

    # Slide 7: Key Discovery #1
    add_content_slide(
        prs,
        "Key Discovery #1: Strategic Silence Dominates",
        [
            "RL learned minimal intervention across all regimes:",
            "",
            "• Very High Stakes (ratio=15): 0.00 interruptions",
            "• Balanced (ratio=3): 0.47 interruptions",
            "• Moderate Low (ratio=2): 3.37 interruptions",
            "",
            "Insight: Even with costly failures, staying silent is often optimal",
            "when baseline failure risk is high (f0=0.6)"
        ]
    )

    # Slide 8: Key Discovery #2
    add_content_slide(
        prs,
        "Key Discovery #2: Inverse Cost Relationship",
        [
            "Counterintuitive finding: Lower stakes → MORE interventions",
            "",
            "Expected (Hypothesis 1):",
            "  • Very high stakes → proactive intervention",
            "  • Moderate low → strategic silence",
            "",
            "Actual (OPPOSITE!):",
            "  • Very high stakes → complete silence (0 interruptions)",
            "  • Moderate low → selective intervention (3.37 interruptions)",
            "",
            "Why? Interruption costs dominate when failure prevention is difficult"
        ]
    )

    # Slide 9: Key Discovery #3
    add_content_slide(
        prs,
        "Key Discovery #3: Simple Tasks Excel",
        [
            "Contrary to expectations, simple tasks show largest gains:",
            "",
            "Simple tasks (8-9 steps): 30-36% improvement",
            "  • make_cereal: 35.77%",
            "  • make_tea: 33.08%",
            "  • make_coffee: 32.66%",
            "  • make_sandwich: 30.99%",
            "",
            "Complex tasks (17-20 steps): 5-19% improvement",
            "  • make_stencil: 5.60%",
            "  • latte_making: 12.51%",
            "",
            "RL optimization is most effective on straightforward procedures"
        ]
    )

    # Slide 10: Performance Heatmap
    fig_path = PROJECT_ROOT / "results" / "figures" / "cross_task_performance_heatmap.png"
    add_image_slide(
        prs,
        "Performance Heatmap: Task × Regime",
        fig_path,
        "RL mean reward and improvement % for all 21 task-regime combinations"
    )

    # Slide 11: Policy Comparison
    fig_path = PROJECT_ROOT / "results" / "figures" / "cross_task_policy_comparison.png"
    add_image_slide(
        prs,
        "Policy Comparison Across Regimes",
        fig_path,
        "RL_PPO consistently outperforms Random, Proactive, and Reactive baselines"
    )

    # Slide 12: Interruption Patterns
    fig_path = PROJECT_ROOT / "results" / "figures" / "cross_task_interruptions_heatmap.png"
    add_image_slide(
        prs,
        "Intervention Frequency Patterns",
        fig_path,
        "RL learns context-dependent intervention strategies (note inverse cost relationship)"
    )

    # Slide 13: Failure Analysis
    fig_path = PROJECT_ROOT / "results" / "figures" / "cross_task_failures_heatmap.png"
    add_image_slide(
        prs,
        "Failure Rate Analysis",
        fig_path,
        "RL accepts some failures to avoid costly interruptions"
    )

    # Slide 14: Summary Dashboard
    fig_path = PROJECT_ROOT / "results" / "figures" / "cross_task_summary_dashboard.png"
    add_image_slide(
        prs,
        "Comprehensive Results Dashboard",
        fig_path,
        "Multi-panel view of performance, improvements, interventions, and tradeoffs"
    )

    # Slide 15: Research Implications
    add_content_slide(
        prs,
        "Research Implications",
        [
            "1. Strategic silence is robust across tasks and cost structures",
            "",
            "2. Interruption costs dominate decision-making even in high-stakes",
            "",
            "3. Simple tasks are low-hanging fruit for RL optimization (30-36% gains)",
            "",
            "4. Complex tasks need special treatment (curriculum learning, longer training)",
            "",
            "5. Same RL algorithm generalizes across diverse procedural domains",
            "    (cooking, technical, crafting)"
        ]
    )

    # Slide 16: Design Principles
    add_content_slide(
        prs,
        "Design Principles for AI Assistants",
        [
            "From these experiments, we derive principles for real-world systems:",
            "",
            "• Infer interruption cost from context (user focus, task phase)",
            "",
            "• Only interrupt when expected benefit > interruption cost",
            "",
            "• Track implicit memory, remind only when significantly decayed",
            "",
            "• When uncertain, bias toward silence (graceful degradation)",
            "",
            "• Let users control cost parameters ('How disruptive? 1-10')"
        ]
    )

    # Slide 17: Key Insight
    add_content_slide(
        prs,
        "Central Insight",
        [
            '"To be helpful, an assistant must sometimes do nothing."',
            "",
            "This research challenges the assumption that more assistance",
            "is always better.",
            "",
            "By explicitly modeling interruption costs, we reveal that:",
            "  • Context determines optimal strategy",
            "  • Same RL algorithm adapts to different cost structures",
            "  • Sophisticated assistance requires knowing when NOT to help",
            "",
            "Applicable to healthcare, manufacturing, education, and beyond."
        ]
    )

    # Slide 18: Future Work
    add_content_slide(
        prs,
        "Future Research Directions",
        [
            "1. Vary failure risk (f0 ∈ {0.3, 0.4, 0.5, 0.6}) to find strategy switching points",
            "",
            "2. Longer training for complex tasks (200k timesteps)",
            "",
            "3. Curriculum learning: simple → complex task transfer",
            "",
            "4. Hierarchical RL: decompose tasks into sub-procedures",
            "",
            "5. Multi-task learning: single model for all 7 tasks",
            "",
            "6. Human evaluation: test with real users in AR/VR environments"
        ]
    )

    # Slide 19: Conclusions
    add_content_slide(
        prs,
        "Conclusions",
        [
            "✓ RL generalizes across tasks: cooking, technical, crafting domains",
            "",
            "✓ 100% success rate: All 21 models beat baselines (mean +24.21%)",
            "",
            "✓ Strategic silence is robust: Emerges across cost structures",
            "",
            "⚠ Context-dependent but inverted: Lower stakes → more interventions",
            "",
            "⚠ Simple tasks benefit most: 30-36% vs 5-19% for complex tasks",
            "",
            "Main takeaway: Sophisticated AI assistance requires knowing",
            "when NOT to help. This validates 'strategic silence' across",
            "diverse tasks and cost structures."
        ]
    )

    # Slide 20: Deliverables
    add_content_slide(
        prs,
        "Experiment Deliverables",
        [
            "Code: 3 new Python scripts (~1,500 lines total)",
            "  • train_cross_task_multi_regime.py",
            "  • evaluate_cross_task_multi_regime.py",
            "  • visualize_cross_task_multi_regime.py",
            "",
            "Models: 21 trained RL policies (all available)",
            "",
            "Data: Training logs + evaluation results (517 KB JSON)",
            "",
            "Visualizations: 10 publication-quality figures (PNG + PDF)",
            "",
            "Documentation: Comprehensive analysis document",
            "  • docs/english/CROSS_TASK_MULTI_REGIME_RESULTS.md"
        ]
    )

    # Save presentation
    output_path = PROJECT_ROOT / "results" / "presentations" / "Cross_Task_Multi_Regime_Results.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    print(f"✓ Presentation saved to: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")

    return output_path


def main():
    """Main entry point."""
    print("="*70)
    print("GENERATING CROSS-TASK MULTI-REGIME PRESENTATION")
    print("="*70)
    print()

    output_path = create_presentation()

    print()
    print("="*70)
    print("PRESENTATION COMPLETE")
    print("="*70)
    print(f"PowerPoint file: {output_path}")
    print(f"Size: {output_path.stat().st_size / 1024:.1f} KB")
    print("="*70)


if __name__ == '__main__':
    main()
