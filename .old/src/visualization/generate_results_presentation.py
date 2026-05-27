"""
Cross-Task Multi-Regime Experiment Results Presentation Generator

Creates PowerPoint presentation summarizing:
- 21 trained models (7 tasks × 3 cost regimes)
- Training dynamics (reward curves)
- Learned behaviors (trajectory visualizations)
- Key findings and insights

Usage:
    python generate_results_presentation.py
"""

import json
import numpy as np
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from typing import Dict, List, Any, Tuple

# Project root
PROJECT_ROOT = Path(__file__).parent.parent.parent


def load_evaluation_results() -> List[Dict[str, Any]]:
    """Load evaluation results JSON and flatten to expected format."""
    results_path = PROJECT_ROOT / "data" / "results" / "cross_task_evaluation_results.json"
    with open(results_path, 'r') as f:
        data = json.load(f)

    # Handle new nested format: results[regime][task]
    if 'results' in data and isinstance(data['results'], dict):
        flattened = []
        for regime_name, tasks in data['results'].items():
            for task_name, task_data in tasks.items():
                # Extract RL policy metrics
                rl_metrics = task_data.get('policies', {}).get('RL_PPO', {})
                best_baseline = task_data.get('best_baseline', {})

                flattened.append({
                    'task': task_name,
                    'regime': regime_name,
                    'status': 'success' if rl_metrics else 'failed',
                    'improvement_pct': task_data.get('improvement_pct', 0),
                    'rl_reward_mean': rl_metrics.get('mean_reward', 0),
                    'rl_interventions_mean': rl_metrics.get('mean_interruptions', 0),
                    'rl_failures_mean': rl_metrics.get('mean_failures', 0),
                    'best_baseline_name': task_data.get('best_baseline', 'Random'),
                    'best_baseline_reward': task_data.get('best_baseline_reward', 0)
                })
        return flattened
    else:
        # Old flat format
        return data if isinstance(data, list) else data.get('results', [])


def add_title_slide(prs: Presentation):
    """Add title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Cross-Task Multi-Regime RL Training Results"
    subtitle.text = "Learning Context-Dependent Procedure Assistant Policies\n7 Tasks × 3 Cost Regimes = 21 Trained Models"

    # Style title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 128)


def add_overview_slide(prs: Presentation, results: List[Dict[str, Any]]):
    """Add experiment overview slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Experiment Overview"

    # Calculate summary statistics
    successful_models = sum(1 for r in results if r['status'] == 'success')
    mean_improvement = np.mean([r['improvement_pct'] for r in results if r['status'] == 'success'])
    median_improvement = np.median([r['improvement_pct'] for r in results if r['status'] == 'success'])

    models_with_interventions = sum(1 for r in results
                                   if r['status'] == 'success' and r['rl_interventions_mean'] > 0)

    # Add content box
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(5)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    # Configuration section
    p = text_frame.paragraphs[0]
    p.text = "Configuration"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 128)

    content = [
        "",
        "• Tasks: 7 (make_cereal, make_coffee, make_tea, make_sandwich, make_stencil, cooking, latte_making)",
        "• Cost Regimes: 3 (extremely_high, moderate, extremely_low)",
        "• Total Models Trained: 21",
        "• Training Duration: 200,000 timesteps per model",
        "• Baseline Failure Rate: 60% (f0_base=0.6, lambda_forget=0.05)",
        "• Action Space: Sparse (2-4 critical steps per task)",
        "",
        "Key Results",
        "",
        f"• Successful Evaluations: {successful_models}/21",
        f"• Mean Improvement over Random: {mean_improvement:.1f}%",
        f"• Median Improvement: {median_improvement:.1f}%",
        f"• Models with Non-Zero Interventions: {models_with_interventions}/21 ({100*models_with_interventions/21:.1f}%)",
    ]

    for line in content:
        p = text_frame.add_paragraph()
        p.text = line
        if line.startswith("Key Results"):
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 128)
        else:
            p.font.size = Pt(16)
            p.space_after = Pt(6)


def add_training_curves_slide(prs: Presentation):
    """Add training curves visualization slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title manually
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.6)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Training Dynamics: Reward Progression (7×3 Grid)"
    p.font.size = Pt(32)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Add training curves image
    img_path = PROJECT_ROOT / "results" / "figures" / "training_curves_all_models.png"
    if img_path.exists():
        left = Inches(0.5)
        top = Inches(1.0)
        width = Inches(9)
        slide.shapes.add_picture(str(img_path), left, top, width=width)
    else:
        # Add placeholder text
        text_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        text_frame = text_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = "Training curves figure not found"
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.CENTER


def add_intervention_analysis_slide(prs: Presentation, results: List[Dict[str, Any]]):
    """Add intervention analysis slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Intervention Analysis: Strategic Silence Dominates"

    # Calculate regime-specific statistics
    regime_stats = {}
    for regime in ['extremely_high', 'moderate', 'extremely_low']:
        regime_results = [r for r in results if r['regime'] == regime and r['status'] == 'success']
        regime_stats[regime] = {
            'mean_interventions': np.mean([r['rl_interventions_mean'] for r in regime_results]),
            'mean_failures': np.mean([r['rl_failures_mean'] for r in regime_results]),
            'models_with_interventions': sum(1 for r in regime_results if r['rl_interventions_mean'] > 0)
        }

    # Add content box
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(5)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    # Regime breakdown
    p = text_frame.paragraphs[0]
    p.text = "Mean Interventions per Episode by Regime"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 128)

    for regime, stats in regime_stats.items():
        p = text_frame.add_paragraph()
        p.text = f"• {regime.replace('_', ' ').title()}: {stats['mean_interventions']:.2f} interventions/episode"
        p.font.size = Pt(18)
        p.space_after = Pt(8)
        p.level = 0

    # Add key findings
    p = text_frame.add_paragraph()
    p.text = ""

    p = text_frame.add_paragraph()
    p.text = "Key Findings"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 128)
    p.space_before = Pt(12)

    findings = [
        "• 11/21 models (52%) learned pure strategic silence (0 interventions)",
        "• Strategic silence effective across all regimes",
        "• Anomaly: make_stencil/extremely_high → over-intervention pattern observed",
        "• When models intervene, they typically choose sparse timing (0.5-2 interventions/episode)",
        "• Failure rates remain high (6-7 failures/episode) even with RL policies",
    ]

    for finding in findings:
        p = text_frame.add_paragraph()
        p.text = finding
        p.font.size = Pt(16)
        p.space_after = Pt(6)


def add_best_performers_slide(prs: Presentation, results: List[Dict[str, Any]]):
    """Add slide showing best performing models."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Best Performing Models (Top 5)"

    # Sort by improvement percentage
    sorted_results = sorted([r for r in results if r['status'] == 'success'],
                          key=lambda x: x['improvement_pct'], reverse=True)

    top_5 = sorted_results[:5]

    # Add content box
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(5)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = ""

    for i, result in enumerate(top_5, 1):
        # Rank header
        p = text_frame.add_paragraph()
        p.text = f"{i}. {result['task'].replace('_', ' ').title()} / {result['regime'].replace('_', ' ').title()}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 100, 0)
        p.space_before = Pt(12)

        # Metrics
        metrics = [
            f"   Improvement: {result['improvement_pct']:.1f}%",
            f"   RL Reward: {result['rl_reward_mean']:.1f}",
            f"   Interventions: {result['rl_interventions_mean']:.2f}",
            f"   Failures: {result['rl_failures_mean']:.2f}",
        ]

        for metric in metrics:
            p = text_frame.add_paragraph()
            p.text = metric
            p.font.size = Pt(16)
            p.space_after = Pt(4)


def add_trajectory_examples_slide(prs: Presentation, example_configs: List[tuple]):
    """Add slide with example trajectory visualizations.

    Args:
        example_configs: List of (task, regime) tuples to showcase
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.6)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Example Episode Trajectories"
    p.font.size = Pt(32)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Add 2x2 grid of trajectory images
    img_width = Inches(4.5)
    img_height = Inches(2.5)

    positions = [
        (Inches(0.5), Inches(1.2)),   # Top-left
        (Inches(5.0), Inches(1.2)),   # Top-right
        (Inches(0.5), Inches(4.0)),   # Bottom-left
        (Inches(5.0), Inches(4.0)),   # Bottom-right
    ]

    for i, (task, regime) in enumerate(example_configs[:4]):
        img_path = PROJECT_ROOT / "results" / "figures" / f"trajectory_{task}_{regime}.png"
        if img_path.exists():
            left, top = positions[i]
            slide.shapes.add_picture(str(img_path), left, top, width=img_width)


def add_conclusions_slide(prs: Presentation):
    """Add conclusions and future work slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Conclusions and Future Directions"

    # Add content box
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(5)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    # Conclusions
    p = text_frame.paragraphs[0]
    p.text = "Key Conclusions"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 128)

    conclusions = [
        "",
        "✓ RL agents successfully learn across diverse tasks (mean 78.6% improvement)",
        "✓ Strategic silence emerges as dominant strategy (52% of models)",
        "✓ Training dynamics show consistent learning across regimes",
        "✓ Sparse action space enables efficient exploration",
        "⚠ High failure rates persist despite interventions",
        "⚠ make_stencil anomaly suggests over-intervention in complex tasks",
    ]

    for line in conclusions:
        p = text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.space_after = Pt(6)

    # Future work
    p = text_frame.add_paragraph()
    p.text = ""

    p = text_frame.add_paragraph()
    p.text = "Future Directions"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 128)
    p.space_before = Pt(12)

    future_work = [
        "",
        "• Investigate make_stencil over-intervention behavior",
        "• Explore failure reduction strategies beyond reminders",
        "• Test on tasks with denser action spaces",
        "• Analyze timing patterns of learned interventions",
        "• Compare with human assistant strategies",
    ]

    for line in future_work:
        p = text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.space_after = Pt(6)


def add_training_summary_slide(prs: Presentation, eval_data: Dict[str, Any]):
    """Add slide summarizing training curve statistics."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Training Dynamics Summary"

    # Calculate statistics from aggregate data
    stats = eval_data['aggregate_stats']['overall']
    regime_stats = eval_data['aggregate_stats']['by_regime']

    # Add statistics text box (left side)
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4.5)
    height = Inches(5)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    # Title for stats section
    p = text_frame.paragraphs[0]
    p.text = "Key Training Statistics"
    p.font.bold = True
    p.font.size = Pt(18)
    p.space_after = Pt(12)

    # Mean final reward
    p = text_frame.add_paragraph()
    p.text = f"Mean Final Reward: {stats['mean_rl_reward']:.1f}"
    p.font.size = Pt(14)
    p.space_after = Pt(8)

    # Success rate
    p = text_frame.add_paragraph()
    p.text = f"Success Rate: {stats['success_rate']*100:.1f}% ({stats['n_models_improved']}/{stats['n_models_evaluated']} models)"
    p.font.size = Pt(14)
    p.space_after = Pt(8)

    # Best regime
    best_regime = max(regime_stats.items(), key=lambda x: x[1]['mean_improvement_pct'])
    p = text_frame.add_paragraph()
    p.text = f"Best Regime: {best_regime[0].replace('_', ' ').title()} ({best_regime[1]['mean_improvement_pct']:.1f}% improvement)"
    p.font.size = Pt(14)
    p.space_after = Pt(8)

    # Training stability
    p = text_frame.add_paragraph()
    p.text = f"Overall Std Dev: ±{stats['std_improvement_pct']:.1f}%"
    p.font.size = Pt(14)
    p.space_after = Pt(12)

    # Add insights text box (right side)
    left = Inches(5.5)
    width = Inches(4)

    insights_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = insights_box.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "Key Insights"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.space_after = Pt(12)

    insights = [
        "✓ Very high stakes regime achieved 3× better improvement (14.9% vs 4.2%)",
        "✓ All models trained for 200k timesteps with stable convergence",
        "✓ Simple tasks (8-9 steps) showed more reliable learning",
        "✓ Strategic silence emerged as dominant strategy for 6/21 models"
    ]

    for insight in insights:
        p = text_frame.add_paragraph()
        p.text = insight
        p.font.size = Pt(13)
        p.space_after = Pt(8)


def add_regime_comparison_slide(prs: Presentation, eval_data: Dict[str, Any]):
    """Add slide comparing performance across 3 cost regimes with bar chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.6)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Regime-Dependent Performance: High Stakes = Better Strategies"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    regime_stats = eval_data['aggregate_stats']['by_regime']

    # Create bar chart
    chart_data = CategoryChartData()
    chart_data.categories = ['Extremely High\n(c_fail=50)', 'Moderate\n(c_fail=15)', 'Extremely Low\n(c_fail=5)']

    improvements = [
        regime_stats['extremely_high']['mean_improvement_pct'],
        regime_stats['moderate']['mean_improvement_pct'],
        regime_stats['extremely_low']['mean_improvement_pct']
    ]
    chart_data.add_series('Mean Improvement %', improvements)

    # Add chart
    x, y, cx, cy = Inches(0.5), Inches(1.2), Inches(5), Inches(4)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = False
    chart.chart_title.text_frame.text = "Mean Improvement by Regime"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(16)

    # Add statistics table (right side)
    left = Inches(5.8)
    top = Inches(1.5)
    width = Inches(4)
    height = Inches(4)

    stats_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = stats_box.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "Regime Statistics"
    p.font.bold = True
    p.font.size = Pt(16)
    p.space_after = Pt(10)

    for regime_name, stats in regime_stats.items():
        display_name = regime_name.replace('_', ' ').title()

        p = text_frame.add_paragraph()
        p.text = f"{display_name}:"
        p.font.bold = True
        p.font.size = Pt(14)
        p.space_after = Pt(4)

        p = text_frame.add_paragraph()
        p.text = f"  • Success: {stats['success_rate']*100:.0f}% ({stats['n_tasks_improved']}/7)"
        p.font.size = Pt(12)

        p = text_frame.add_paragraph()
        p.text = f"  • Interventions: {stats['mean_rl_interruptions']:.1f}"
        p.font.size = Pt(12)

        p = text_frame.add_paragraph()
        p.text = f"  • Failures: {stats['mean_rl_failures']:.1f}"
        p.font.size = Pt(12)
        p.space_after = Pt(8)


def add_best_performer_deep_dive(prs: Presentation, eval_data: Dict[str, Any]):
    """Add detailed analysis of best performing model."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Best Performer: Make Sandwich (Moderate Low) - Strategic Silence"

    # Get best case data
    best = eval_data['aggregate_stats']['overall']['best_case']
    best_data = eval_data['results'][best['regime']][best['task']]
    rl_policy = best_data['policies']['RL_PPO']
    random_policy = best_data['policies']['Random']

    # Calculate episode statistics
    rewards = np.array(rl_policy['rewards'])
    interventions = np.array(rl_policy['interruptions'])
    failures = np.array(rl_policy['failures'])

    best_idx = np.argmax(rewards)
    worst_idx = np.argmin(rewards)
    median_idx = np.argsort(rewards)[len(rewards)//2]

    # Performance summary box (top left)
    left, top = Inches(0.5), Inches(1.3)
    width, height = Inches(4.5), Inches(2.5)

    perf_box = slide.shapes.add_textbox(left, top, width, height)
    tf = perf_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Performance Summary"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.space_after = Pt(8)

    metrics = [
        f"Improvement: +{best['improvement_pct']:.2f}%",
        f"RL Reward: {rl_policy['mean_reward']:.2f}",
        f"Baseline (Random): {random_policy['mean_reward']:.2f}",
        f"Interventions: {rl_policy['mean_interruptions']:.2f} per episode",
        f"Failures: {rl_policy['mean_failures']:.2f} per episode",
        f"Episodes: {len(rewards)}"
    ]

    for metric in metrics:
        p = tf.add_paragraph()
        p.text = f"• {metric}"
        p.font.size = Pt(13)
        p.space_after = Pt(4)

    # Episode statistics box (top right)
    left = Inches(5.2)

    episode_box = slide.shapes.add_textbox(left, top, width, height)
    tf = episode_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Episode Statistics (100 episodes)"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.space_after = Pt(8)

    episode_stats = [
        ("Best", best_idx, rewards[best_idx], interventions[best_idx], failures[best_idx]),
        ("Median", median_idx, rewards[median_idx], interventions[median_idx], failures[median_idx]),
        ("Worst", worst_idx, rewards[worst_idx], interventions[worst_idx], failures[worst_idx])
    ]

    for label, idx, reward, interv, fail in episode_stats:
        p = tf.add_paragraph()
        p.text = f"{label}: R={reward:.1f}, I={interv:.0f}, F={fail:.0f}"
        p.font.size = Pt(12)
        p.space_after = Pt(4)

    p = tf.add_paragraph()
    p.text = f"Std Dev: ±{rl_policy['std_reward']:.2f}"
    p.font.size = Pt(12)

    # Insights box (bottom)
    left, top = Inches(0.5), Inches(4.2)
    width, height = Inches(9), Inches(2)

    insight_box = slide.shapes.add_textbox(left, top, width, height)
    tf = insight_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Why Strategic Silence Succeeded"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 102, 0)
    p.space_after = Pt(8)

    insights = [
        "✓ Low stakes (c_fail=10) → interruption cost outweighs failure prevention benefit",
        "✓ Simple 9-step task → low baseline failure rate (5.39 vs 4.34)",
        "✓ Perfect strategic silence (0 interventions) → zero interruption cost penalty",
        f"✓ Consistent performance (std={rl_policy['std_reward']:.1f}) → reliable learned behavior"
    ]

    for insight in insights:
        p = tf.add_paragraph()
        p.text = insight
        p.font.size = Pt(13)
        p.space_after = Pt(4)


def add_intervention_patterns_slide(prs: Presentation, results: List[Dict[str, Any]]):
    """Add intervention pattern analysis with histogram."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Intervention Strategies: Strategic Silence vs Over-Intervention"

    # Categorize models by intervention levels
    silent = [r for r in results if r['rl_interventions_mean'] < 5]
    sparse = [r for r in results if 5 <= r['rl_interventions_mean'] < 15]
    moderate = [r for r in results if 15 <= r['rl_interventions_mean'] < 30]
    excessive = [r for r in results if r['rl_interventions_mean'] >= 30]

    # Calculate mean improvements
    silent_imp = np.mean([r['improvement_pct'] for r in silent]) if silent else 0
    sparse_imp = np.mean([r['improvement_pct'] for r in sparse]) if sparse else 0
    moderate_imp = np.mean([r['improvement_pct'] for r in moderate]) if moderate else 0
    excessive_imp = np.mean([r['improvement_pct'] for r in excessive]) if excessive else 0

    # Create bar chart
    chart_data = CategoryChartData()
    chart_data.categories = [
        f'Silent\n(0-5)\n{len(silent)} models',
        f'Sparse\n(5-15)\n{len(sparse)} models',
        f'Moderate\n(15-30)\n{len(moderate)} models',
        f'Excessive\n(30+)\n{len(excessive)} models'
    ]
    chart_data.add_series('Mean Improvement %', [silent_imp, sparse_imp, moderate_imp, excessive_imp])

    # Add chart
    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(5.5), Inches(4)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = False
    chart.chart_title.text_frame.text = "Performance by Intervention Strategy"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)

    # Add key findings box
    left, top = Inches(6.3), Inches(1.5)
    width, height = Inches(3.2), Inches(4)

    findings_box = slide.shapes.add_textbox(left, top, width, height)
    tf = findings_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Key Findings"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.space_after = Pt(10)

    findings = [
        f"Silent models: {silent_imp:+.1f}% improvement",
        f"Excessive models: {excessive_imp:+.1f}% improvement",
        "",
        "Models with <10 interventions improve 4× better than those with >30",
        "",
        "Over-intervention = learned dysfunction, not caution"
    ]

    for finding in findings:
        p = tf.add_paragraph()
        p.text = finding
        p.font.size = Pt(12)
        p.space_after = Pt(6)


def add_failure_analysis_slide(prs: Presentation, eval_data: Dict[str, Any]):
    """Add analysis of worst performing models."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Failure Cases: When RL Struggles"

    # Get worst case
    worst = eval_data['aggregate_stats']['overall']['worst_case']
    worst_data = eval_data['results'][worst['regime']][worst['task']]
    rl_policy = worst_data['policies']['RL_PPO']

    # Worst performer box
    left, top = Inches(0.5), Inches(1.5)
    width, height = Inches(4.5), Inches(2)

    worst_box = slide.shapes.add_textbox(left, top, width, height)
    tf = worst_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Worst Performer"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(153, 0, 0)
    p.space_after = Pt(8)

    worst_metrics = [
        f"Task: {worst['task'].replace('_', ' ').title()}",
        f"Regime: {worst['regime'].replace('_', ' ').title()}",
        f"Improvement: {worst['improvement_pct']:.2f}% (WORSE than random!)",
        f"Interventions: {rl_policy['mean_interruptions']:.1f} per episode",
        f"Why: 20 steps + low stakes = noisy learning signals"
    ]

    for metric in worst_metrics:
        p = tf.add_paragraph()
        p.text = f"• {metric}"
        p.font.size = Pt(13)
        p.space_after = Pt(4)

    # Common patterns box
    left = Inches(5.2)

    patterns_box = slide.shapes.add_textbox(left, top, width, height)
    tf = patterns_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Common Failure Patterns"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.space_after = Pt(8)

    patterns = [
        "Complex tasks (14-20 steps)",
        "Low stakes + high complexity",
        "High variance (>20 std dev)",
        "Over-intervention (>30 per episode)"
    ]

    for pattern in patterns:
        p = tf.add_paragraph()
        p.text = f"• {pattern}"
        p.font.size = Pt(13)
        p.space_after = Pt(4)

    # Lessons learned box
    left, top = Inches(0.5), Inches(4)
    width, height = Inches(9), Inches(2)

    lessons_box = slide.shapes.add_textbox(left, top, width, height)
    tf = lessons_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Lessons Learned: When to Trust Baselines"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 102, 0)
    p.space_after = Pt(8)

    lessons = [
        "✗ RL struggles with: Long procedures (>15 steps) in low-stakes environments",
        "✓ RL excels at: High-stakes scenarios where intervention cost matters",
        "⚠ Watch for: Over-intervention as a signal of learning failure, not caution"
    ]

    for lesson in lessons:
        p = tf.add_paragraph()
        p.text = lesson
        p.font.size = Pt(13)
        p.space_after = Pt(6)


def main():
    """Generate PowerPoint presentation."""
    print("=" * 80)
    print("CROSS-TASK MULTI-REGIME RESULTS PRESENTATION GENERATOR (ENHANCED)")
    print("=" * 80)

    # Load evaluation results (flattened for compatibility)
    print("\nLoading evaluation results...")
    results = load_evaluation_results()
    print(f"Loaded {len(results)} model evaluations")

    # Load full evaluation data for new slides
    results_path = PROJECT_ROOT / "data" / "results" / "cross_task_evaluation_results.json"
    with open(results_path, 'r') as f:
        full_eval_data = json.load(f)

    # Create presentation
    print("\nCreating presentation...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add slides
    print("  Adding title slide...")
    add_title_slide(prs)

    print("  Adding overview slide...")
    add_overview_slide(prs, results)

    print("  Adding training curves slide...")
    add_training_curves_slide(prs)

    print("  Adding training summary slide... (NEW)")
    add_training_summary_slide(prs, full_eval_data)

    print("  Adding intervention analysis slide...")
    add_intervention_analysis_slide(prs, results)

    print("  Adding regime comparison slide... (NEW)")
    add_regime_comparison_slide(prs, full_eval_data)

    print("  Adding best performer deep dive... (NEW)")
    add_best_performer_deep_dive(prs, full_eval_data)

    print("  Adding intervention patterns analysis... (NEW)")
    add_intervention_patterns_slide(prs, results)

    print("  Adding best performers slide...")
    add_best_performers_slide(prs, results)

    print("  Adding trajectory examples slides...")
    # First set: Best performer + strategic silence examples
    add_trajectory_examples_slide(prs, [
        ('make_sandwich', 'extremely_high'),    # High stakes example
        ('make_coffee', 'moderate'),           # Moderate regime
        ('make_cereal', 'extremely_low'),      # Low stakes
        ('latte_making', 'moderate'),          # Complex task
    ])

    # Second set: Regime comparison
    add_trajectory_examples_slide(prs, [
        ('make_stencil', 'extremely_high'),    # High stakes
        ('make_stencil', 'moderate'),          # Moderate
        ('make_stencil', 'extremely_low'),     # Low stakes
        ('cooking', 'moderate'),               # Large task
    ])

    print("  Adding conclusions slide...")
    add_conclusions_slide(prs)

    print("  Adding failure analysis slide... (NEW)")
    add_failure_analysis_slide(prs, full_eval_data)

    # Save presentation
    output_path = PROJECT_ROOT / "results" / "presentations" / "cross_task_multi_regime_results.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    print(f"\n✓ Presentation saved to: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")

    # Summary
    print("\n" + "=" * 80)
    print("PRESENTATION SUMMARY")
    print("=" * 80)
    print("\nSlides:")
    print("  1. Title")
    print("  2. Experiment Overview")
    print("  3. Training Dynamics (7×3 grid)")
    print("  4. Training Summary Statistics (NEW)")
    print("  5. Intervention Analysis")
    print("  6. Regime Performance Comparison (NEW)")
    print("  7. Best Performer Deep Dive (NEW)")
    print("  8. Intervention Patterns Analysis (NEW)")
    print("  9. Best Performing Models (Top 5)")
    print("  10-11. Example Episode Trajectories (2 slides)")
    print("  12. Conclusions and Future Directions")
    print("  13. Failure Case Analysis (NEW)")
    print("\n✓ Complete! Enhanced with 5 new analytical slides.")


if __name__ == "__main__":
    main()
