#!/usr/bin/env python3
"""Generate research presentation for MA-IPPO Procedure Assistant project.

Slides:
 1. Title
 2. Problem Setting
 3. Dec-POMDP System Overview (diagram)
 4. Agent Spaces: Observation & Actions
 5. Belief Update Mechanism
 6. Simulation Parameters (4 regime axes)
 7. Failure Model
 8. Tasks Overview
 9. IPPO Training Procedure
10. Results: Cross-Task Summary
11-17. Learning Curves (7 tasks)
18. Key Findings
19. Discussion
20. Next Steps

Usage:
    python generate_research_presentation.py [--output path.pptx]
"""
import io
import json
import glob
from pathlib import Path

import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.patheffects as pe
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PROJECT_ROOT = Path(__file__).parent.parent.parent

# ── Palette ────────────────────────────────────────────────────────────────
C_BG       = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK     = RGBColor(0x1A, 0x1A, 0x2E)
C_ACCENT   = RGBColor(0x16, 0x21, 0x3E)
C_BLUE     = RGBColor(0x0F, 0x3F, 0x6E)
C_TEAL     = RGBColor(0x0F, 0x6E, 0x6E)
C_RED      = RGBColor(0xD3, 0x2F, 0x2F)
C_GREEN    = RGBColor(0x38, 0x8E, 0x3C)
C_ORANGE   = RGBColor(0xF5, 0x7C, 0x00)
C_GRAY     = RGBColor(0x60, 0x60, 0x60)
C_LGRAY    = RGBColor(0xF5, 0xF5, 0xF5)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

REGIME_COLORS_MPL = {
    'extremely_low':  '#1976D2',
    'balanced':       '#388E3C',
    'extremely_high': '#D32F2F',
}

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)


# ═══════════════════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════════════════

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_layout(prs):
    return prs.slide_layouts[6]   # completely blank


def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_width_pt=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
                wrap=True, italic=False):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or C_DARK
    return txb


def add_image(slide, img_path, left, top, width, height=None):
    if height:
        slide.shapes.add_picture(str(img_path),
                                 Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    else:
        slide.shapes.add_picture(str(img_path),
                                 Inches(left), Inches(top),
                                 Inches(width))


def header_bar(slide, title, subtitle=None):
    """Dark header bar at top of slide."""
    add_rect(slide, 0, 0, 13.33, 1.10, fill_rgb=C_DARK)
    add_textbox(slide, title, 0.3, 0.10, 12.0, 0.65,
                font_size=24, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle, 0.3, 0.72, 12.0, 0.35,
                    font_size=13, color=RGBColor(0xBB, 0xBB, 0xBB),
                    align=PP_ALIGN.LEFT)


def bullet_block(slide, items, left, top, width, height,
                 font_size=16, color=None, indent_by_level=True):
    """Add a multi-level bullet list. items = list of (level, text) tuples."""
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for level, text in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = level
        indent = "    " * level if indent_by_level else ""
        bullet_char = "•" if level == 0 else "–" if level == 1 else "·"
        run = p.add_run()
        run.text = f"{indent}{bullet_char}  {text}"
        run.font.size = Pt(font_size - level * 1.5)
        run.font.color.rgb = color or C_DARK


def fig_to_png_bytes(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    buf.seek(0)
    return buf


def add_fig(slide, fig, left, top, width, height=None):
    buf = fig_to_png_bytes(fig)
    if height:
        slide.shapes.add_picture(buf, Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    else:
        slide.shapes.add_picture(buf, Inches(left), Inches(top),
                                 Inches(width))
    plt.close(fig)


# ═══════════════════════════════════════════════════════════════════════════
# DIAGRAM GENERATORS
# ═══════════════════════════════════════════════════════════════════════════

def make_pomdp_diagram():
    """Dec-POMDP system overview as matplotlib figure."""
    fig, ax = plt.subplots(figsize=(13, 6))
    ax.set_xlim(0, 13)
    ax.set_ylim(0, 6)
    ax.axis('off')

    def box(cx, cy, w, h, label, sublabel=None, fc='#FFFFFF', ec='#333333', lw=1.5, fs=11, sfs=9):
        rect = FancyBboxPatch((cx - w/2, cy - h/2), w, h,
                              boxstyle="round,pad=0.08",
                              facecolor=fc, edgecolor=ec, linewidth=lw, zorder=3)
        ax.add_patch(rect)
        y_text = cy + (0.15 if sublabel else 0)
        ax.text(cx, y_text, label, ha='center', va='center',
                fontsize=fs, fontweight='bold', zorder=4)
        if sublabel:
            ax.text(cx, cy - 0.28, sublabel, ha='center', va='center',
                    fontsize=sfs, color='#555555', zorder=4)

    def arrow(x1, y1, x2, y2, label='', color='#333333', style='->', lw=1.5, dashed=False):
        ls = '--' if dashed else '-'
        ax.annotate('', xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle=style, color=color,
                                   lw=lw, linestyle=ls),
                    zorder=2)
        if label:
            mx, my = (x1 + x2) / 2, (y1 + y2) / 2
            ax.text(mx, my + 0.15, label, ha='center', va='bottom',
                    fontsize=8, color=color, zorder=5,
                    bbox=dict(fc='white', ec='none', pad=1))

    # ── Background group boxes ──────────────────────────────────────────
    # Hidden state
    bg_state = FancyBboxPatch((0.2, 4.1), 12.9, 1.6,
                               boxstyle="round,pad=0.05",
                               facecolor='#F8F8FF', edgecolor='#AAAAAA',
                               linewidth=1, linestyle='--', zorder=1)
    ax.add_patch(bg_state)
    ax.text(0.5, 5.55, "LATENT STATE  sₜ  (hidden from assistant)",
            fontsize=9, color='#888888')

    # Human agent group
    bg_human = FancyBboxPatch((0.2, 0.3), 5.2, 3.5,
                               boxstyle="round,pad=0.05",
                               facecolor='#EEF4FF', edgecolor='#1976D2',
                               linewidth=1.5, zorder=1)
    ax.add_patch(bg_human)
    ax.text(1.0, 3.65, "HUMAN  (full-info MDP)",
            fontsize=10, fontweight='bold', color='#1976D2')

    # Assistant agent group
    bg_asst = FancyBboxPatch((7.6, 0.3), 5.5, 3.5,
                              boxstyle="round,pad=0.05",
                              facecolor='#FFF3E0', edgecolor='#F57C00',
                              linewidth=1.5, zorder=1)
    ax.add_patch(bg_asst)
    ax.text(9.2, 3.65, "ASSISTANT  (POMDP, belief-based)",
            fontsize=10, fontweight='bold', color='#F57C00')

    # ── State nodes ─────────────────────────────────────────────────────
    box(2.0, 5.0, 2.6, 0.8, "sₜ  =  (step, τ, memory[N], obs_noise)",
        fc='#E8EAF6', ec='#3F51B5', lw=2, fs=10)
    box(8.5, 5.0, 2.2, 0.8, "f(m) = f₀·e^{−k·m}",
        sublabel="failure model",
        fc='#FFEBEE', ec='#D32F2F', lw=1.5, fs=10)

    # ── Human nodes ─────────────────────────────────────────────────────
    box(1.8, 2.8, 3.6, 0.9, "oᴴ  (dim=5)",
        sublabel="step · τ · memory · asst_action · obs_noise",
        fc='#BBDEFB', ec='#1976D2', lw=1.5, fs=10)
    box(1.8, 1.6, 3.0, 0.7, "πᴴ  —  PPO MLP 64×64",
        fc='#E3F2FD', ec='#1976D2', lw=1.5, fs=10)
    box(1.8, 0.7, 3.8, 0.7, "Aᴴ  :  silent / narrate / question_j",
        fc='#BBDEFB', ec='#1976D2', lw=1.5, fs=10)

    # ── Assistant nodes ─────────────────────────────────────────────────
    box(10.5, 2.8, 3.2, 0.9, "bₜ  =  P(step=s)  [N+1]",
        sublabel="Bayesian belief via noisy obs + prior",
        fc='#FFE0B2', ec='#F57C00', lw=2, fs=10)
    box(10.5, 1.6, 3.0, 0.7, "πᴬ  —  PPO MLP 64×64",
        fc='#FFF8E1', ec='#F57C00', lw=1.5, fs=10)
    box(10.5, 0.7, 3.8, 0.7, "Aᴬ  :  silent / confirm / remind_j",
        fc='#FFE0B2', ec='#F57C00', lw=1.5, fs=10)

    # ── Shared reward ────────────────────────────────────────────────────
    box(6.5, 0.7, 2.6, 0.7, "Shared Reward  rₜ",
        sublabel="−c_fail  −c_nar  −c_q  −c_remind  −c_confirm  +R_complete",
        fc='#E8F5E9', ec='#388E3C', lw=2, fs=10)

    # ── Arrows ───────────────────────────────────────────────────────────
    # State → human obs (perfect)
    arrow(2.0, 4.6, 1.8, 3.25, "perfect", color='#1976D2', lw=2)
    # State → assistant belief (noisy)
    arrow(3.2, 4.6, 9.5, 3.25, "noisy (obs_noise)", color='#F57C00', lw=2, dashed=True)
    # State → failure → reward
    arrow(8.5, 4.6, 8.5, 4.15)
    arrow(8.5, 4.15, 6.8, 1.05, "c_fail", color='#D32F2F', lw=1.5)

    # Human: obs→policy→action
    arrow(1.8, 2.35, 1.8, 1.95)
    arrow(1.8, 1.25, 1.8, 1.05)

    # Assistant: belief→policy→action
    arrow(10.5, 2.35, 10.5, 1.95)
    arrow(10.5, 1.25, 10.5, 1.05)

    # NARRATION: human action → belief hard-reset
    ax.annotate('', xy=(9.0, 0.85), xytext=(3.7, 0.70),
                arrowprops=dict(arrowstyle='->', color='#D32F2F',
                                lw=2.5, connectionstyle='arc3,rad=-0.35'),
                zorder=2)
    ax.text(6.5, 0.3, "narrate  →  belief hard-reset", ha='center',
            fontsize=9, color='#D32F2F', fontweight='bold', zorder=5)

    # Human action → state (question: memory up, narrate: noise down)
    arrow(2.8, 0.7, 3.5, 4.6, "memory↑ / noise↓", color='#555555', lw=1.2, dashed=True)
    # Assistant action → state (remind: memory up)
    arrow(9.0, 0.7, 8.0, 4.6, "memory↑", color='#555555', lw=1.2, dashed=True)

    # Actions → reward
    arrow(1.8, 0.35, 5.5, 0.7, "c_nar, c_q", color='#388E3C', lw=1.2)
    arrow(10.5, 0.35, 7.5, 0.7, "c_remind, c_confirm", color='#388E3C', lw=1.2)

    # Last-action visibility (cross-agent)
    ax.annotate('', xy=(9.3, 2.8), xytext=(3.6, 2.8),
                arrowprops=dict(arrowstyle='<->', color='#777777',
                                lw=1, linestyle='--'),
                zorder=2)
    ax.text(6.5, 2.95, "last action visible", ha='center',
            fontsize=8, color='#777777')

    fig.tight_layout(pad=0.2)
    return fig


def make_belief_update_diagram():
    """Belief update flow diagram."""
    fig, ax = plt.subplots(figsize=(12, 5))
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 5)
    ax.axis('off')

    def box(cx, cy, w, h, lines, fc='#FFFFFF', ec='#333333', lw=1.5, fs=10):
        rect = FancyBboxPatch((cx - w/2, cy - h/2), w, h,
                              boxstyle="round,pad=0.10",
                              facecolor=fc, edgecolor=ec, linewidth=lw, zorder=3)
        ax.add_patch(rect)
        n = len(lines)
        for i, (txt, style) in enumerate(lines):
            dy = (i - (n-1)/2) * 0.32
            ax.text(cx, cy - dy, txt, ha='center', va='center',
                    fontsize=style.get('fs', fs),
                    fontweight=style.get('fw', 'normal'),
                    color=style.get('color', '#111111'),
                    family=style.get('family', 'sans-serif'),
                    zorder=4)

    def arr(x1, y1, x2, y2, label='', color='#333333', lw=1.5, dashed=False, rad=0):
        ls = '--' if dashed else '-'
        cs = f'arc3,rad={rad}' if rad else 'arc3,rad=0'
        ax.annotate('', xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle='->', color=color, lw=lw,
                                   linestyle=ls, connectionstyle=cs), zorder=2)
        if label:
            mx, my = (x1+x2)/2, (y1+y2)/2
            ax.text(mx+0.05, my+0.18, label, ha='center', fontsize=8, color=color)

    # ── Tick start ──────────────────────────────────────────────────────
    box(2.0, 4.2, 3.4, 0.9,
        [("TICK t  —  simultaneous actions", {'fw':'bold', 'fs':10}),
         ("bₜ₋₁ : prior belief    aᴴ, aᴬ : joint actions", {'fs':8.5, 'color':'#444'})],
        fc='#E8EAF6', ec='#3F51B5', lw=2)

    # ── Decision diamond ────────────────────────────────────────────────
    diamond = plt.Polygon([[3.8, 2.9],[4.6, 2.5],[3.8, 2.1],[3.0, 2.5]],
                           closed=True, fc='#FFF3E0', ec='#F57C00', lw=2, zorder=3)
    ax.add_patch(diamond)
    ax.text(3.8, 2.5, "narrate?", ha='center', va='center',
            fontsize=9, fontweight='bold', color='#F57C00', zorder=4)

    arr(2.0, 3.75, 2.0, 2.75)
    arr(2.0, 2.75, 3.0, 2.5)

    # ── PATH A: NO narration ─────────────────────────────────────────────
    ax.text(2.5, 1.55, "NO", fontsize=9, color='#555', ha='center')
    arr(3.0, 2.5, 2.0, 2.0, rad=0.0)

    box(1.0, 1.35, 3.0, 1.0,
        [("1. Prior propagation", {'fw':'bold', 'fs':9.5}),
         ("b̃[s] = (1−h)·b[s] + h·(Tᵀb)[s]", {'fs':8.5, 'family':'monospace', 'color':'#1A237E'}),
         ("h = 1/30  (step advance prob.)", {'fs':8, 'color':'#555'})],
        fc='#E3F2FD', ec='#1976D2', lw=1.5)

    box(4.5, 1.35, 3.2, 1.0,
        [("2. Bayesian update", {'fw':'bold', 'fs':9.5}),
         ("ℓ[s] = (1−p)·𝟙(s=obs) + p·𝒩(obs−s)", {'fs':8.5, 'family':'monospace', 'color':'#1A237E'}),
         ("bₜ ∝ b̃ · ℓ    where p = obs_noise_state", {'fs':8, 'color':'#555'})],
        fc='#E3F2FD', ec='#1976D2', lw=1.5)

    arr(1.0, 0.85, 4.5, 0.85)

    # ── Noise recovery ───────────────────────────────────────────────────
    box(2.8, 0.40, 3.8, 0.55,
        [("obs_noise recovers: n(t) = baseline−(baseline−min)·e^{−λt}", {'fs':8.5, 'color':'#555'})],
        fc='#F5F5F5', ec='#AAAAAA', lw=1)
    arr(4.5, 0.85, 3.8, 0.68, dashed=True)

    # ── PATH B: YES narration ─────────────────────────────────────────────
    ax.text(5.5, 2.65, "YES", fontsize=9, color='#D32F2F', ha='center')
    arr(4.6, 2.5, 7.5, 2.5, color='#D32F2F', lw=2)

    box(8.5, 2.5, 3.0, 1.1,
        [("Hard reset (perfect info)", {'fw':'bold', 'fs':9.5, 'color':'#B71C1C'}),
         ("bₜ[true_step] = 1.0", {'fs':9, 'family':'monospace', 'color':'#C62828'}),
         ("obs_noise_state ← obs_noise_min", {'fs':8.5, 'family':'monospace', 'color':'#555'}),
         ("effect persists: half-life = ln2/λ ticks", {'fs':8, 'color':'#777'})],
        fc='#FFEBEE', ec='#D32F2F', lw=2)

    # ── Output ───────────────────────────────────────────────────────────
    box(8.5, 0.8, 2.8, 0.8,
        [("bₜ  →  assistant policy input", {'fw':'bold', 'fs':10}),
         ("step_belief[N+1]  posterior", {'fs':8.5, 'color':'#555'})],
        fc='#E8F5E9', ec='#388E3C', lw=2)

    arr(6.0, 1.35, 7.1, 0.95)
    arr(8.5, 1.95, 8.5, 1.20)

    fig.tight_layout(pad=0.2)
    return fig


def make_ippo_diagram():
    """IPPO alternating training diagram."""
    fig, ax = plt.subplots(figsize=(10, 4.5))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 4.5)
    ax.axis('off')

    def box(cx, cy, w, h, text, fc='#FFFFFF', ec='#333333', lw=1.5, fs=11, fw='normal'):
        rect = FancyBboxPatch((cx - w/2, cy - h/2), w, h,
                              boxstyle="round,pad=0.10",
                              facecolor=fc, edgecolor=ec, linewidth=lw, zorder=3)
        ax.add_patch(rect)
        ax.text(cx, cy, text, ha='center', va='center',
                fontsize=fs, fontweight=fw, zorder=4, wrap=True,
                multialignment='center')

    def arr(x1, y1, x2, y2, label='', color='#333', lw=1.5):
        ax.annotate('', xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle='->', color=color, lw=lw), zorder=2)
        if label:
            ax.text((x1+x2)/2, (y1+y2)/2+0.15, label,
                    ha='center', fontsize=8.5, color=color)

    # Initialize
    box(5, 4.1, 3.5, 0.65, "Initialize πᴴ, πᴬ  (random)", fc='#E8EAF6', ec='#3F51B5', lw=2, fw='bold')

    # Loop
    box(2.0, 2.8, 3.0, 0.75, "Train  πᴴ  (50k steps)\nπᴬ frozen", fc='#BBDEFB', ec='#1976D2', lw=1.5)
    box(8.0, 2.8, 3.0, 0.75, "Train  πᴬ  (50k steps)\nπᴴ frozen", fc='#FFE0B2', ec='#F57C00', lw=1.5)
    box(5.0, 1.7, 3.5, 0.75, "Joint Eval  (200 episodes)\nmean_reward  /  mean_failures", fc='#F3E5F5', ec='#7B1FA2', lw=1.5)
    box(5.0, 0.6, 3.5, 0.65, "Save best checkpoint\nif reward improved", fc='#E8F5E9', ec='#388E3C', lw=1.5)

    arr(5, 3.78, 5, 3.35)
    arr(3.5, 3.35, 2.0, 3.18)
    arr(2.0, 2.43, 2.0, 2.0)
    arr(2.0, 2.0, 3.5, 1.70, "→ round")
    arr(6.5, 1.70, 8.0, 2.0)
    arr(8.0, 2.43, 8.0, 2.0)
    arr(8.0, 2.0, 6.5, 1.70)
    arr(5.0, 1.33, 5.0, 0.93)

    # Repeat arrow
    ax.annotate('', xy=(3.5, 3.18), xytext=(5.0, 0.28),
                arrowprops=dict(arrowstyle='->', color='#777', lw=1.2,
                                connectionstyle='arc3,rad=0.5'),
                zorder=2)
    ax.text(1.0, 1.8, "repeat\n15 rounds", ha='center', fontsize=9, color='#777')

    fig.tight_layout(pad=0.2)
    return fig


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    # gradient-like background using two rectangles
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=C_DARK)
    add_rect(slide, 0, 4.8, 13.33, 2.7, fill_rgb=RGBColor(0x0D, 0x47, 0xA1))

    add_textbox(slide,
                "Learning Context-Dependent\nProcedure Assistant Policies",
                1.0, 1.0, 11.0, 2.8,
                font_size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                "Dec-POMDP  +  IPPO  —  Multi-Agent RL for Cooperative Assistance",
                1.0, 3.8, 11.0, 0.6,
                font_size=18, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)

    add_textbox(slide,
                "7 tasks  ×  3 failure-cost regimes  ×  15 training rounds\n"
                "comm=cheap_narrate  |  decay=step_transition  |  obs=durable",
                1.0, 5.2, 11.0, 1.2,
                font_size=14, color=RGBColor(0x90, 0xCA, 0xF9), align=PP_ALIGN.CENTER)
    return slide


def slide_problem(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, "Problem Setting",
               "When should an AI assistant speak up — and when stay silent?")

    items = [
        (0, "Human performs sequential procedural tasks (cooking, crafting, …)"),
        (1, "N steps, each with a specific order; some steps are safety-critical"),
        (1, "Execution order varies across rollouts (partial flexibility)"),
        (0, "Human has memory that degrades between steps"),
        (1, "Memory determines failure probability  f(m) = f₀ · exp(−k · m)"),
        (1, "Without support: f₀ = 60% failure at zero memory"),
        (0, "AI assistant observes (partially) and must decide each tick:"),
        (1, "Remind a critical step?  Confirm what the human is doing?  Stay silent?"),
        (0, "Core tension:"),
        (1, "Interrupting too much  →  interruption cost, slows human"),
        (1, "Interrupting too little  →  steps fail, high failure cost"),
        (0, "Research question: Can RL learn context-dependent policies that adapt to cost structure?"),
    ]
    bullet_block(slide, items, 0.5, 1.25, 12.5, 6.0, font_size=16)
    return slide


def slide_pomdp(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, "Dec-POMDP System Overview",
               "Human: full-info MDP  |  Assistant: POMDP with Bayesian belief  |  Shared cooperative reward")
    fig = make_pomdp_diagram()
    add_fig(slide, fig, 0.15, 1.15, 13.0, 6.2)
    return slide


def slide_agent_spaces(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, "Agent Observation & Action Spaces",
               "Key asymmetry: Human has perfect state info; Assistant infers via Bayesian belief")

    # Left column: Human
    add_rect(slide, 0.3, 1.2, 5.9, 5.8, fill_rgb=RGBColor(0xEE, 0xF4, 0xFF),
             line_rgb=RGBColor(0x19, 0x76, 0xD2), line_width_pt=1.5)
    add_textbox(slide, "HUMAN AGENT  (MDP)", 0.5, 1.25, 5.5, 0.45,
                font_size=13, bold=True, color=RGBColor(0x0D, 0x47, 0xA1))

    human_obs = [
        (0, "Observation  oᴴ  (dim = 5)"),
        (1, "current_step  —  exact (perfect channel)"),
        (1, "τ  —  elapsed ticks (exact)"),
        (1, "memory[current_step]  (own memory, exact)"),
        (1, "assistant_last_action"),
        (1, "obs_noise_state"),
        (0, "Actions  Aᴴ  (size = 2 + Nc)"),
        (1, "silent"),
        (1, "narrate  →  belief hard-reset  +  obs_noise ↓"),
        (1, "question_j  →  memory[crit_j] += 0.4"),
    ]
    bullet_block(slide, human_obs, 0.5, 1.75, 5.6, 5.0, font_size=13)

    # Right column: Assistant
    add_rect(slide, 7.1, 1.2, 5.9, 5.8, fill_rgb=RGBColor(0xFF, 0xF3, 0xE0),
             line_rgb=RGBColor(0xF5, 0x7C, 0x00), line_width_pt=1.5)
    add_textbox(slide, "ASSISTANT AGENT  (POMDP)", 7.3, 1.25, 5.5, 0.45,
                font_size=13, bold=True, color=RGBColor(0xE6, 0x51, 0x00))

    asst_obs = [
        (0, "Observation  oᴬ  (dim = N+1 + Nc + 1)"),
        (1, "step_belief[N+1]  —  Bayesian posterior P(step=s)"),
        (1, "memory_estimate_critical[Nc]  (inferred, never ground truth)"),
        (1, "human_last_action"),
        (0, "Actions  Aᴬ  (size = 2 + Nc)"),
        (1, "silent"),
        (1, "confirm  →  costs c_confirm"),
        (1, "remind_j  →  memory[crit_j] += 0.8"),
        (2, "+ c_off_timing penalty if wrong step"),
    ]
    bullet_block(slide, asst_obs, 7.3, 1.75, 5.6, 5.0, font_size=13)

    # Divider
    add_textbox(slide, "↔  last action visible  ↔", 6.05, 3.5, 1.25, 0.5,
                font_size=10, color=C_GRAY, align=PP_ALIGN.CENTER)
    return slide


def slide_belief_update(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, "Belief Update Mechanism  (per tick)",
               "Prior propagation + Bayesian update each tick  |  Narration bypasses with hard reset")
    fig = make_belief_update_diagram()
    add_fig(slide, fig, 0.2, 1.15, 12.9, 6.1)
    return slide


def slide_parameters(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, "Simulation Parameters  —  4 Orthogonal Axes",
               "Each experiment = one combination of (fail_regime, comm_regime, decay_regime, obs_regime)")

    # 4 columns for 4 axes
    col_configs = [
        ("1. Failure Cost",  0.3,  RGBColor(0xD3, 0x2F, 0x2F), [
            "extremely_low:   scale = 2.0",
            "balanced:          scale = 15.0",
            "extremely_high: scale = 50.0",
            "",
            "c_fail_per_step =",
            "  scale × criticality",
        ]),
        ("2. Comm Cost",  3.6, RGBColor(0x19, 0x76, 0xD2), [
            "default:",
            "  c_nar=0.5  c_remind=1.0",
            "  c_q=0.5    c_confirm=1.0",
            "cheap_narrate:",
            "  c_nar=0.1  c_remind=1.0",
            "  c_q=0.5    c_confirm=0.5",
            "c_off_timing = 3.0 (both)",
        ]),
        ("3. Memory Decay",  6.9, RGBColor(0x38, 0x8E, 0x3C), [
            "default:",
            "  λ_forget = 0.03",
            "  memory_init = 0.0",
            "step_transition:",
            "  λ_forget = 0.10",
            "  memory_init = 0.3",
            "",
            "Decay at step transition:",
            "  m *= (1−λ)^elapsed",
        ]),
        ("4. Obs Noise",  10.2, RGBColor(0xF5, 0x7C, 0x00), [
            "default:",
            "  λ_recover = 0.10",
            "  half-life ≈ 7 ticks",
            "  (< 1 step avg.)",
            "durable:",
            "  λ_recover = 0.02",
            "  half-life ≈ 35 ticks",
            "  (≈ 1 step avg.)",
            "obs_noise_min = 0.05",
        ]),
    ]

    for title, left, color, lines in col_configs:
        add_rect(slide, left, 1.15, 2.9, 5.95,
                 fill_rgb=RGBColor(0xFA, 0xFA, 0xFA),
                 line_rgb=color, line_width_pt=2)
        add_textbox(slide, title, left + 0.1, 1.20, 2.7, 0.45,
                    font_size=13, bold=True, color=color)
        content = "\n".join(lines)
        add_textbox(slide, content, left + 0.15, 1.70, 2.7, 5.2,
                    font_size=11, color=C_DARK)

    # Fixed params footer
    add_rect(slide, 0.3, 7.1, 12.7, 0.32, fill_rgb=RGBColor(0xF0, 0xF0, 0xF0))
    add_textbox(slide,
                "Fixed: f₀=0.6  k=3.0  δ_remind=0.8  δ_q=0.4  R_complete=10.0  "
                "β=1.0  step_dur=30±10 tick  |  PPO: lr=3e-4  n_steps=2048  batch=64  epochs=10  γ=0.99",
                0.4, 7.12, 12.5, 0.28, font_size=9, color=C_GRAY)
    return slide


def slide_failure_model(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, "Failure Model  &  Task Overview",
               "Discrete hazard model with memory-dependent failure probability")

    # Failure model plot
    fig, ax = plt.subplots(figsize=(5, 3.5))
    m = np.linspace(0, 1.5, 200)
    for k_label, k_val, color in [("k=3.0 (current)", 3.0, '#D32F2F'),
                                    ("k=1.5", 1.5, '#1976D2')]:
        f = 0.6 * np.exp(-k_val * m)
        ax.plot(m, f, color=color, lw=2, label=k_label)
    ax.axhline(0, color='#999', lw=0.5)
    ax.set_xlabel("memory  m", fontsize=11)
    ax.set_ylabel("failure prob  f(m)", fontsize=11)
    ax.set_title("f(m) = 0.6 · exp(−k · m)", fontsize=12, fontweight='bold')
    ax.legend(fontsize=9)
    ax.set_xlim(0, 1.5)
    ax.set_ylim(0, 0.65)
    ax.grid(True, alpha=0.3)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    fig.tight_layout()
    add_fig(slide, fig, 0.3, 1.2, 5.2, 3.8)

    # Task table
    tasks = [
        ("make_coffee",   8, 1, 12.0, "cooking"),
        ("make_cereal",   8, 2, 10.0, "cooking"),
        ("make_tea",      9, 2, 12.0, "cooking"),
        ("make_sandwich", 9, 1, 15.0, "cooking"),
        ("cooking",      14, 4, 20.0, "cooking"),
        ("make_stencil", 17, 4, 30.0, "crafting"),
        ("latte_making", 20, 2, 25.0, "technical"),
    ]
    headers = ["Task", "N", "Nc", "c_fail base", "Domain"]
    col_widths = [2.3, 0.6, 0.6, 1.1, 1.3]
    col_lefts = [5.7, 8.0, 8.6, 9.2, 10.3]
    row_h = 0.42

    add_textbox(slide, "7 Tasks", 5.7, 1.2, 6.2, 0.38,
                font_size=14, bold=True, color=C_DARK)

    # Header row
    add_rect(slide, 5.6, 1.6, 7.4, row_h, fill_rgb=C_DARK)
    for hdr, lft, w in zip(headers, col_lefts, col_widths):
        add_textbox(slide, hdr, lft, 1.62, w, row_h - 0.04,
                    font_size=11, bold=True, color=C_WHITE)

    # Data rows
    for i, (name, n, nc, cf, dom) in enumerate(tasks):
        y = 1.6 + (i + 1) * row_h
        bg = RGBColor(0xF5, 0xF5, 0xF5) if i % 2 == 0 else C_WHITE
        add_rect(slide, 5.6, y, 7.4, row_h, fill_rgb=bg,
                 line_rgb=RGBColor(0xDD, 0xDD, 0xDD), line_width_pt=0.5)
        for val, lft, w in zip([name, n, nc, cf, dom], col_lefts, col_widths):
            add_textbox(slide, str(val), lft, y + 0.02, w, row_h - 0.04, font_size=11)

    return slide


def slide_ippo(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, "IPPO Training Procedure",
               "Alternating best-response: fix one agent, train the other; repeat for 15 rounds")

    fig = make_ippo_diagram()
    add_fig(slide, fig, 0.5, 1.2, 7.0, 4.5)

    items = [
        (0, "Independent PPO (IPPO)"),
        (1, "Each agent has its own policy network"),
        (1, "No shared parameters between agents"),
        (0, "Alternating best-response training"),
        (1, "Odd round: fix πᴬ, train πᴴ for 50k steps"),
        (1, "Even round: fix πᴴ, train πᴬ for 50k steps"),
        (0, "Checkpoint saving"),
        (1, "Save if joint eval reward improves"),
        (1, "Final eval uses best checkpoint"),
        (0, "Total compute per task × regime"),
        (1, "15 rounds × 2 × 50k = 1.5M env steps"),
        (1, "3 regimes = 4.5M steps per task"),
    ]
    bullet_block(slide, items, 7.8, 1.3, 5.2, 6.0, font_size=14)
    return slide


def slide_cross_task(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, "Results: Cross-Task Summary",
               "Final eval across 7 tasks × 3 failure-cost regimes  |  comm=cheap_narrate  decay=step_transition  obs=durable")
    fig_path = PROJECT_ROOT / "results" / "figures" / "cross_task_durable.png"
    if fig_path.exists():
        add_image(slide, fig_path, 0.3, 1.15, 12.7, 6.1)
    return slide


def slide_learning_curves(prs, task_name, json_path):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, f"Learning Curves: {task_name.replace('_', ' ').title()}",
               "Reward, Narration, Questions, Remind+Confirm, Failures  vs  training round")
    fig_path = PROJECT_ROOT / "results" / "figures" / f"learning_curves_{task_name}_durable.png"
    if fig_path.exists():
        add_image(slide, fig_path, 0.2, 1.15, 12.9, 5.5)

    # Load summary numbers
    try:
        d = json.load(open(json_path))
        lines = []
        for regime, rv in d.get('regimes', {}).items():
            fe = rv.get('ma_ippo', {}).get('final_eval', {})
            r = fe.get('mean_reward', 0)
            f = fe.get('mean_failures', 0)
            n = fe.get('mean_narrations', 0)
            lines.append(f"{regime:18s}  reward={r:+.2f}  fail={f:.2f}  narr={n:.2f}")
        add_textbox(slide, "\n".join(lines), 0.3, 6.75, 12.5, 0.65,
                    font_size=10, color=C_GRAY)
    except Exception:
        pass
    return slide


def slide_findings(prs, all_results):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, "Key Findings",
               "What the agents learned across 7 tasks × 3 failure-cost regimes")

    items = [
        (0, "Strategic silence is the dominant strategy at low failure cost"),
        (1, "Agents learn to trust human and stay silent → minimal interruption cost"),
        (0, "Narration is rarely learned (obs=durable reduces spurious narration loops)"),
        (1, "Most narration appears only in high-cost regimes with complex tasks"),
        (1, "make_stencil (17 steps) still shows narr≈59 at extremely_high"),
        (0, "Question-asking is the primary cooperation signal at medium/high cost"),
        (1, "make_cereal balanced: narr=17 — cereal has 2 critical steps, questioning pays off"),
        (0, "Remind-based strategy emerges for assistant when failure cost is high"),
        (1, "cooking extremely_high: interact=6.77/ep, failures reduced to 0.28"),
        (0, "Task complexity drives strategy divergence"),
        (1, "Simple tasks (make_coffee, make_sandwich): reward stays positive even at high cost"),
        (1, "Complex tasks (latte_making 20 steps): harder convergence, reward −52 at high cost"),
        (0, "Durable obs regime (half-life 35 ticks) successfully eliminates narration=100 loops"),
        (1, "make_stencil: was narr=101 with default obs → now narr=59 (1 narrate/critical step)"),
    ]
    bullet_block(slide, items, 0.4, 1.2, 12.5, 6.0, font_size=14)
    return slide


def slide_discussion(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, "Discussion",
               "Interpretation, limitations, and open questions")

    items = [
        (0, "Why does IPPO converge to strategic silence in low-cost regimes?"),
        (1, "Interruption costs outweigh expected failure cost savings → optimal to stay quiet"),
        (1, "Both agents independently learn this equilibrium without communication"),
        (0, "Narration vs. questioning: different information channels"),
        (1, "Narration: hard-resets assistant belief, but costs c_nar each time"),
        (1, "Question: boosts human memory directly, cheaper with cheap_narrate regime"),
        (0, "The durable obs regime changes narration economics"),
        (1, "Longer-lasting effect per narration → fewer narrations needed"),
        (1, "Agents find narration more 'efficient' when each narration persists longer"),
        (0, "Limitations"),
        (1, "IPPO may not converge to true Nash equilibrium — only local best responses"),
        (1, "15 rounds may not be sufficient for complex tasks (latte_making)"),
        (1, "Simulation assumes discrete ticks; real tasks are continuous"),
        (1, "Human policy is also learned — may not reflect real human behavior"),
        (0, "The 4-axis regime structure enables systematic ablation"),
        (1, "Each axis can be varied independently to isolate effects"),
    ]
    bullet_block(slide, items, 0.4, 1.2, 12.5, 6.0, font_size=14)
    return slide


def slide_phase_diagram(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, "Phase Diagram: Emergent Communication Regimes",
               "15×15 grid search  |  c_comm (=c_nar=c_q) × c_fail_scale  |  c_remind=1.0 fixed  |  make_cereal")
    fig_path = PROJECT_ROOT / "results" / "figures" / "phase_diagram_15x15.png"
    if fig_path.exists():
        add_image(slide, fig_path, 0.2, 1.15, 12.9, 4.2)

    items = [
        (0, "Clear phase transitions emerge along two axes:"),
        (1, "Low c_fail (< 5): Silent — interruption cost dominates, agents learn to stay quiet"),
        (1, "High c_fail + low c_comm: Question-dominated — cheap info-seeking to prevent costly failures"),
        (1, "High c_fail + high c_comm: Remind-only — assistant uses minimal targeted interventions"),
        (0, "Boundary at c_fail ≈ 5–15 marks the transition from silence to active cooperation"),
        (0, "Narration is rarely used — question is strictly preferred when c_q = c_nar"),
    ]
    bullet_block(slide, items, 0.3, 5.5, 12.7, 2.0, font_size=13)
    return slide


def slide_next_steps(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, "Next Steps",
               "Future work directions")

    items = [
        (0, "Extend phase diagram analysis"),
        (1, "Run grid search on complex tasks (cooking, make_stencil, latte_making)"),
        (1, "Add asymmetric cost axis: vary c_nar independently from c_q"),
        (1, "Higher resolution at phase boundaries for sharper transitions"),
        (0, "Ablation studies across all 4 regime axes"),
        (1, "Compare durable vs default obs regime across all tasks"),
        (1, "Isolate memory decay effect: default vs step_transition"),
        (0, "More training rounds for complex tasks"),
        (1, "latte_making extremely_high shows −52 reward — needs more rounds or larger network"),
        (1, "Try 30 rounds or adaptive stopping based on convergence"),
        (0, "Richer agent architectures"),
        (1, "Replace MLP with LSTM to handle temporal dependencies within a step"),
        (1, "Attention over belief state for assistant"),
        (0, "Real human evaluation"),
        (1, "Replace learned human policy with human subjects in a user study"),
        (1, "Measure whether learned assistant policy transfers to real humans"),
        (0, "Additional task domains"),
        (1, "Medical procedures, lab protocols, manufacturing checklists"),
        (1, "Tasks with irreversible failures (higher stakes than c_fail=50)"),
    ]
    bullet_block(slide, items, 0.4, 1.2, 12.5, 6.0, font_size=14)
    return slide


# ═══════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════

def create_presentation(output_path=None):
    if output_path is None:
        output_path = str(PROJECT_ROOT / "results" / "presentations" /
                          "research_presentation_durable.pptx")

    prs = new_prs()

    slide_title(prs)
    slide_problem(prs)
    slide_pomdp(prs)
    slide_agent_spaces(prs)
    slide_belief_update(prs)
    slide_parameters(prs)
    slide_failure_model(prs)
    slide_ippo(prs)
    slide_cross_task(prs)

    # Learning curves — 7 tasks
    all_results = {}
    task_order = ["make_coffee", "make_cereal", "make_tea", "make_sandwich",
                  "cooking", "make_stencil", "latte_making"]
    for task in task_order:
        json_path = (PROJECT_ROOT / "data" / "results" /
                     f"ma_{task}_cheap_narrate_step_transition_durable.json")
        if json_path.exists():
            d = json.load(open(json_path))
            all_results[task] = d
            slide_learning_curves(prs, task, json_path)

    slide_findings(prs, all_results)
    slide_phase_diagram(prs)
    slide_discussion(prs)
    slide_next_steps(prs)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"Saved: {output_path}  ({len(prs.slides)} slides)")


if __name__ == '__main__':
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument('--output', default=None)
    args = parser.parse_args()
    create_presentation(args.output)
